#!/usr/bin/env python3
"""Generate Climate Risk Platform overview PPT (13 slides, Korean, sales/customer deck).

Usage:
    source backend/venv/bin/activate
    python generate_ppt.py
"""

import sys
import os
import platform

# Add backend to path so we can import services
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "backend"))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from app.data.sample_facilities import get_all_facilities
from app.services.transition_risk import analyse_scenario, compare_scenarios
from app.services.esg_compliance import assess_framework
from app.services.physical_risk import assess_physical_risk
from app.core.config import SCENARIOS, PROJECTION_YEARS

# ── Color Palette ────────────────────────────────────────────────────
NAVY = RGBColor(0x1E, 0x29, 0x3B)
DARK_NAVY = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
ALT_ROW_GRAY = RGBColor(0xF8, 0xFA, 0xFC)
MID_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY = RGBColor(0x47, 0x55, 0x69)
BLACK = RGBColor(0x1E, 0x1E, 0x1E)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
YELLOW = RGBColor(0xEA, 0xB3, 0x08)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
LIGHT_GREEN = RGBColor(0xDC, 0xFC, 0xE7)
LIGHT_RED = RGBColor(0xFE, 0xE2, 0xE2)
LIGHT_ORANGE = RGBColor(0xFF, 0xED, 0xD5)
LIGHT_YELLOW = RGBColor(0xFE, 0xF9, 0xC3)

SCENARIO_COLORS = {
    "net_zero_2050": RED,
    "below_2c": ORANGE,
    "delayed_transition": YELLOW,
    "current_policies": GREEN,
}
SCENARIO_BG_COLORS = {
    "net_zero_2050": LIGHT_RED,
    "below_2c": LIGHT_ORANGE,
    "delayed_transition": LIGHT_YELLOW,
    "current_policies": LIGHT_GREEN,
}

# ── Slide dimensions (16:9) ─────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Layout constants ─────────────────────────────────────────────────
LEFT_MARGIN = Inches(0.8)
RIGHT_MARGIN = Inches(0.8)
CONTENT_WIDTH = Inches(11.7)

# ── Font strategy ────────────────────────────────────────────────────
_IS_MAC = platform.system() == "Darwin"
LATIN_FONT = "Calibri"
EA_FONT = "Apple SD Gothic Neo" if _IS_MAC else "Malgun Gothic"

# ── Sector name map ──────────────────────────────────────────────────
SECTOR_NAMES_KR = {
    "steel": "철강",
    "petrochemical": "석유화학",
    "automotive": "자동차",
    "electronics": "반도체/전자",
    "utilities": "발전",
    "cement": "시멘트",
    "shipping": "해운",
    "oil_gas": "정유",
}

HAZARD_NAMES_KR = {
    "flood": "홍수",
    "typhoon": "태풍",
    "heatwave": "폭염",
    "drought": "가뭄",
    "sea_level_rise": "해수면",
}


# ══════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ══════════════════════════════════════════════════════════════════════

def _set_font(run, name=LATIN_FONT, ea_name=EA_FONT, size=18, bold=False, color=BLACK):
    """Set font with EA (East Asian) fallback via lxml."""
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # Set EA typeface via XML
    rPr = run._r.get_or_add_rPr()
    # Remove existing ea element if any
    for ea in rPr.findall(qn('a:ea')):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', ea_name)


def _add_shadow(shape):
    """Add drop shadow to a shape via lxml."""
    spPr = shape._element.spPr
    # Remove existing effectLst if any
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    shadow = etree.SubElement(effectLst, qn('a:outerShdw'))
    shadow.set('blurRad', '50800')
    shadow.set('dist', '25400')
    shadow.set('dir', '5400000')
    shadow.set('algn', 'bl')
    shadow.set('rotWithShape', '0')
    srgb = etree.SubElement(shadow, qn('a:srgbClr'))
    srgb.set('val', '000000')
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', '20000')


def _set_slide_bg(slide, color):
    """Set a solid-color slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_gradient_bg(slide, color1_hex, color2_hex, angle=5400000):
    """Set gradient background on a slide via lxml.
    angle: EMU angle. 5400000 = top-to-bottom (270deg).
    """
    bg = slide.background
    bgPr = bg._element
    # Find or create bgPr
    cSld = slide._element
    bg_elem = cSld.find(qn('p:bg'))
    if bg_elem is None:
        bg_elem = etree.SubElement(cSld, qn('p:bg'))
        cSld.insert(0, bg_elem)

    # Remove existing bgPr
    for old in bg_elem.findall(qn('p:bgPr')):
        bg_elem.remove(old)

    bgPr = etree.SubElement(bg_elem, qn('p:bgPr'))
    gradFill = etree.SubElement(bgPr, qn('a:gradFill'))
    gradFill.set('flip', 'none')
    gradFill.set('rotWithShape', '1')

    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))
    # Stop 1
    gs1 = etree.SubElement(gsLst, qn('a:gs'))
    gs1.set('pos', '0')
    srgb1 = etree.SubElement(gs1, qn('a:srgbClr'))
    srgb1.set('val', color1_hex)
    # Stop 2
    gs2 = etree.SubElement(gsLst, qn('a:gs'))
    gs2.set('pos', '100000')
    srgb2 = etree.SubElement(gs2, qn('a:srgbClr'))
    srgb2.set('val', color2_hex)

    lin = etree.SubElement(gradFill, qn('a:lin'))
    lin.set('ang', str(angle))
    lin.set('scaled', '1')

    # effectLst required
    etree.SubElement(bgPr, qn('a:effectLst'))


def _add_rounded_rect(slide, left, top, width, height, fill_color, shadow=True):
    """Add a rounded rectangle with optional drop shadow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if shadow:
        _add_shadow(shape)
    return shape


def _add_shape(slide, left, top, width, height, fill_color):
    """Add a plain rectangle shape (for bars, lines, accents)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a textbox with cross-platform font settings."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, size=font_size, color=color, bold=bold)
    return txBox


def _add_rich_text(slide, left, top, width, height, segments, alignment=PP_ALIGN.LEFT):
    """Add textbox with mixed formatting.
    segments = [(text, font_size, color, bold), ...]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    for text, font_size, color, bold in segments:
        run = p.add_run()
        run.text = text
        _set_font(run, size=font_size, color=color, bold=bold)
    return txBox


def _add_title_bar(slide, title_text, subtitle_text=None):
    """Add a navy title bar at the top of the slide."""
    _add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), NAVY)
    _add_textbox(slide, LEFT_MARGIN, Inches(0.15), Inches(10), Inches(0.6),
                 title_text, font_size=28, color=WHITE, bold=True)
    if subtitle_text:
        _add_textbox(slide, LEFT_MARGIN, Inches(0.7), Inches(10), Inches(0.4),
                     subtitle_text, font_size=14, color=MID_GRAY)
    # Accent line
    _add_shape(slide, Inches(0), Inches(1.2), SLIDE_WIDTH, Inches(0.04), BLUE)


def _add_bullet_list(slide, left, top, width, height, items, font_size=16,
                     color=DARK_GRAY, bullet_char="- ", line_spacing=1.5):
    """Add a bullet list with proper spacing."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"{bullet_char}{item}"
        _set_font(run, size=font_size, color=color)
        p.space_before = Pt(2)
        p.space_after = Pt(int(font_size * (line_spacing - 1)))


def _add_kpi_box(slide, left, top, width, height, label, value,
                 bg_color=LIGHT_BLUE, value_color=BLUE):
    """Add a rounded KPI box with shadow."""
    card = _add_rounded_rect(slide, left, top, width, height, bg_color, shadow=True)
    # Label
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.4),
                 label, font_size=11, color=DARK_GRAY)
    # Value (larger font)
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.45), width - Inches(0.3), Inches(0.5),
                 value, font_size=28, color=value_color, bold=True)


def _add_table(slide, left, top, width, height, headers, rows, col_widths=None,
               header_bg=NAVY, header_fg=WHITE, stripe_color=ALT_ROW_GRAY):
    """Add a styled table to the slide."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        _style_cell(cell, font_size=12, bold=True, color=header_fg, bg_color=header_bg)

    # Rows
    for i, row in enumerate(rows):
        bg = stripe_color if i % 2 == 1 else WHITE
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            if isinstance(val, tuple):
                cell.text = val[0]
                _style_cell(cell, font_size=11, color=val[1],
                            bg_color=bg, bold=val[2] if len(val) > 2 else False)
            else:
                cell.text = str(val)
                _style_cell(cell, font_size=11, color=BLACK, bg_color=bg)

    return table_shape


def _style_cell(cell, font_size=11, bold=False, color=BLACK, bg_color=WHITE):
    """Style a table cell with cross-platform fonts."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg_color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            _set_font(run, size=font_size, bold=bold, color=color)
        if not p.runs:
            run = p.add_run()
            run.text = p.text
            _set_font(run, size=font_size, bold=bold, color=color)
        p.alignment = PP_ALIGN.CENTER


def _add_chevron_flow(slide, steps, y, colors):
    """Add chevron-style process flow diagram."""
    n = len(steps)
    total_width = Inches(12.0)
    gap = Inches(0.1)
    step_width = Emu(int((total_width - gap * (n - 1)) / n))
    start_x = LEFT_MARGIN

    for i, (step_text, color) in enumerate(zip(steps, colors)):
        x = Emu(int(start_x + i * (step_width + gap)))
        shape = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON, x, y, step_width, Inches(0.9)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step_text
        _set_font(run, size=11, color=WHITE, bold=True)


def _add_footer(slide, page_num):
    """Add consistent footer with page number and confidential notice."""
    # Footer line
    _add_shape(slide, LEFT_MARGIN, Inches(7.15), CONTENT_WIDTH, Inches(0.02), MID_GRAY)
    # Page number
    _add_textbox(slide, Inches(12.0), Inches(7.18), Inches(1.0), Inches(0.3),
                 str(page_num), font_size=10, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)
    # Confidential notice
    _add_textbox(slide, LEFT_MARGIN, Inches(7.18), Inches(6.0), Inches(0.3),
                 "Climate Risk Platform  |  Confidential", font_size=10, color=MID_GRAY)


def _add_number_badge(slide, x, y, number_text, bg_color=NAVY, text_color=WHITE, size=Inches(0.45)):
    """Add a circular number badge (e.g., 01, 02, 03)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = number_text
    _set_font(run, size=11, color=text_color, bold=True)
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    return shape


def _fmt_b(val):
    """Format large USD values as $XXX.XB."""
    return f"${abs(val)/1e9:.1f}B"


def _fmt_m(val):
    """Format millions."""
    return f"{val/1e6:.1f}M"


def _fmt_pct(val):
    return f"{val:.0f}%"


# ══════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    """Slide 1: Cover / Title with gradient background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_gradient_bg(slide, '0F172A', '1E293B')

    # Decorative shapes
    _add_shape(slide, Inches(0), Inches(2.8), Inches(0.1), Inches(2.2), BLUE)
    # Subtle decorative rounded rect
    deco = _add_rounded_rect(slide, Inches(10.5), Inches(0.5), Inches(2.5), Inches(2.5),
                             RGBColor(0x1E, 0x29, 0x3B), shadow=False)
    deco.fill.solid()
    deco.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    deco.line.color.rgb = RGBColor(0x2D, 0x3A, 0x50)
    deco.line.width = Pt(1)

    # Title
    _add_textbox(slide, Inches(1.0), Inches(1.8), Inches(11), Inches(1.2),
                 "Climate Risk Platform", font_size=44, color=WHITE, bold=True)
    # Subtitle
    _add_textbox(slide, Inches(1.0), Inches(3.1), Inches(11), Inches(0.7),
                 "NGFS 시나리오 기반 기후 리스크 분석 플랫폼",
                 font_size=22, color=BLUE)
    # Tagline
    _add_textbox(slide, Inches(1.0), Inches(3.9), Inches(11), Inches(0.5),
                 "기후변화 재무 영향 분석  |  ESG 공시 지원 솔루션",
                 font_size=16, color=MID_GRAY)

    # Bottom divider
    _add_shape(slide, Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.02), MID_GRAY)
    _add_textbox(slide, Inches(1.0), Inches(6.15), Inches(11.3), Inches(0.4),
                 "Confidential  |  2026", font_size=12, color=MID_GRAY)


def slide_02_problem(prs):
    """Slide 2: Problems companies face - numbered cards with accent bars."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "기업이 직면한 문제", "기후 리스크 관리, 왜 지금 중요한가?")

    problems = [
        ("ESG 공시 의무화", "2025년부터 ISSB/KSSB 기반 기후 관련 공시 의무 시행\n미준수 시 투자 제한 및 규제 리스크", BLUE),
        ("탄소 비용 급증", "K-ETS 배출권 가격 지속 상승, 직접 재무 영향 확대\nNet Zero 시나리오 기준 2050년 $250/tCO2e", RED),
        ("TCFD 분석 요구", "금융기관/투자자의 기후 리스크 분석 요구 증가\n시나리오 분석, 재무 영향 정량화 필수", ORANGE),
        ("높은 진입장벽", "기후 리스크 분석 = 전문인력 + 복잡한 모델링 필요\n중소/중견기업의 자체 대응 곤란", GREEN),
    ]

    for i, (title, desc, accent_color) in enumerate(problems):
        x = LEFT_MARGIN + (i % 2) * Inches(6.1)
        y = Inches(1.7) + (i // 2) * Inches(2.5)

        # Card background (rounded)
        _add_rounded_rect(slide, x, y, Inches(5.7), Inches(2.1), LIGHT_GRAY, shadow=True)
        # Left accent bar
        _add_shape(slide, x, y + Inches(0.15), Inches(0.06), Inches(1.8), accent_color)

        # Number badge
        _add_number_badge(slide, x + Inches(0.25), y + Inches(0.2),
                          f"{i+1:02d}", bg_color=accent_color)

        # Title
        _add_textbox(slide, x + Inches(0.85), y + Inches(0.2), Inches(4.5), Inches(0.5),
                     title, font_size=18, color=NAVY, bold=True)

        # Description
        _add_textbox(slide, x + Inches(0.35), y + Inches(0.85), Inches(5.0), Inches(1.1),
                     desc, font_size=13, color=DARK_GRAY)

    _add_footer(slide, 2)


def slide_03_service_stages(prs):
    """Slide 3: 3-stage analysis flow (Stage 1→2→3 progressive cards)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "3단계 분석 방법론", "기초 분석에서 캐시플로우 영향까지, 단계적 심층 분석")

    stages = [
        ("Stage 1", "기초 분석", BLUE, LIGHT_BLUE,
         "전환 리스크 개요\n물리적 리스크 개요\n리스크 존재 여부 확인",
         "제공 중"),
        ("Stage 2", "심층 시나리오", ORANGE, LIGHT_ORANGE,
         "NGFS 4개 시나리오 비교\n시설별 NPV 영향 산출\n리스크 등급 순위",
         "제공 중"),
        ("Stage 3", "캐시플로우", GREEN, LIGHT_GREEN,
         "DCF 모델 기반\n전환+물리적 통합 분석\n현금흐름 재무영향 정량화",
         "개발 예정"),
    ]

    card_w = Inches(3.4)
    arrow_w = Inches(0.45)
    gap = Inches(0.15)
    total = 3 * card_w + 2 * (arrow_w + 2 * gap)
    start_x = Emu(int((SLIDE_WIDTH - total) / 2))

    for i, (stage, title, color, bg_color, desc, badge_text) in enumerate(stages):
        x = Emu(int(start_x + i * (card_w + arrow_w + 2 * gap)))
        y = Inches(1.7)

        # Card background
        _add_rounded_rect(slide, x, y, card_w, Inches(3.6), bg_color, shadow=True)
        # Top accent bar
        _add_shape(slide, x, y, card_w, Inches(0.06), color)

        # Stage number badge
        _add_number_badge(slide, Emu(int(x + Inches(0.2))), y + Inches(0.25),
                          str(i + 1), bg_color=color)

        # Stage label
        _add_textbox(slide, Emu(int(x + Inches(0.8))), y + Inches(0.2),
                     Inches(2.4), Inches(0.35),
                     stage, font_size=13, color=color, bold=True)

        # Title
        _add_textbox(slide, Emu(int(x + Inches(0.25))), y + Inches(0.7),
                     Inches(2.9), Inches(0.5),
                     title, font_size=22, color=NAVY, bold=True)

        # Description
        _add_textbox(slide, Emu(int(x + Inches(0.25))), y + Inches(1.35),
                     Inches(2.9), Inches(1.5),
                     desc, font_size=14, color=DARK_GRAY)

        # Badge (제공 중 / 개발 예정)
        badge_bg = color if badge_text == "제공 중" else MID_GRAY
        badge_shape = _add_rounded_rect(
            slide,
            Emu(int(x + Inches(0.25))), y + Inches(3.05),
            Inches(1.3), Inches(0.35), badge_bg, shadow=False,
        )
        tf = badge_shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = badge_text
        _set_font(run, size=11, color=WHITE, bold=True)

        # Arrow between cards (not after last)
        if i < 2:
            ax = Emu(int(x + card_w + gap))
            ay = Inches(3.2)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, ax, ay, arrow_w, Inches(0.5))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MID_GRAY
            arrow.line.fill.background()

    # Bottom navy bar
    _add_rounded_rect(slide, LEFT_MARGIN, Inches(5.8), CONTENT_WIDTH, Inches(0.8), NAVY, shadow=True)
    _add_textbox(slide, LEFT_MARGIN, Inches(5.9), CONTENT_WIDTH, Inches(0.5),
                 "17개 시설 분석  |  NGFS 국제 표준 시나리오  |  K-ETS 한국 정책 반영",
                 font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    _add_footer(slide, 3)


def slide_04_portfolio(prs, facilities):
    """Slide 4: Demo portfolio overview with rounded KPI boxes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "분석 포트폴리오 (데모)", "한국 대표 산업 8개 섹터 커버")

    total_s1 = sum(f["current_emissions_scope1"] for f in facilities)
    total_s2 = sum(f["current_emissions_scope2"] for f in facilities)
    total_assets = sum(f["assets_value"] for f in facilities)
    total_revenue = sum(f["annual_revenue"] for f in facilities)
    companies = set(f["company"] for f in facilities)
    sectors = set(f["sector"] for f in facilities)

    # KPI boxes (rounded + shadow)
    kpis = [
        ("시설 수", f"{len(facilities)}개"),
        ("기업 수", f"{len(companies)}개"),
        ("섹터 수", f"{len(sectors)}개"),
        ("Scope 1+2 배출량", f"{(total_s1+total_s2)/1e6:.0f}M tCO2e"),
        ("총 자산규모", f"${total_assets/1e9:.0f}B"),
        ("총 매출규모", f"${total_revenue/1e9:.0f}B"),
    ]
    for i, (label, value) in enumerate(kpis):
        x = LEFT_MARGIN + (i % 3) * Inches(4.1)
        y = Inches(1.6) + (i // 3) * Inches(1.5)
        _add_kpi_box(slide, x, y, Inches(3.7), Inches(1.1), label, value)

    # Sector breakdown table
    sector_data = {}
    for f in facilities:
        s = f["sector"]
        if s not in sector_data:
            sector_data[s] = {"count": 0, "emissions": 0, "assets": 0}
        sector_data[s]["count"] += 1
        sector_data[s]["emissions"] += f["current_emissions_scope1"] + f["current_emissions_scope2"]
        sector_data[s]["assets"] += f["assets_value"]

    headers = ["섹터", "시설 수", "배출량 (MtCO2e)", "자산 ($B)"]
    rows = []
    for s in ["steel", "utilities", "oil_gas", "petrochemical", "electronics", "cement", "automotive", "shipping"]:
        if s in sector_data:
            d = sector_data[s]
            rows.append([
                SECTOR_NAMES_KR.get(s, s),
                str(d["count"]),
                f"{d['emissions']/1e6:.1f}",
                f"${d['assets']/1e9:.1f}",
            ])

    _add_table(slide, LEFT_MARGIN, Inches(4.7), CONTENT_WIDTH, Inches(2.5),
               headers, rows,
               col_widths=[Inches(3.0), Inches(2.5), Inches(3.1), Inches(3.1)])

    _add_footer(slide, 4)


def slide_05_stage1_overview(prs, comparison, physical_data):
    """Slide 5: Stage 1 basic risk analysis overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Stage 1: 기초 리스크 분석",
                   "전환 리스크와 물리적 리스크, 어떤 위험이 있는가?")

    # ── Left: Transition Risk Summary ──
    _add_rounded_rect(slide, LEFT_MARGIN, Inches(1.6), Inches(5.6), Inches(4.0),
                      LIGHT_BLUE, shadow=True)
    _add_shape(slide, LEFT_MARGIN, Inches(1.6), Inches(5.6), Inches(0.06), BLUE)
    _add_textbox(slide, LEFT_MARGIN + Inches(0.3), Inches(1.8), Inches(5.0), Inches(0.5),
                 "전환 리스크 요약", font_size=20, color=NAVY, bold=True)

    # KPI: Total emissions
    total_emissions = comparison.get("total_portfolio_emissions", 0)
    if total_emissions == 0:
        # Fallback: sum from npv_comparison data
        for npv_item in comparison.get("npv_comparison", []):
            if npv_item.get("scenario") == "current_policies":
                total_emissions = npv_item.get("total_emissions", 0)
                break
    _add_kpi_box(slide, LEFT_MARGIN + Inches(0.3), Inches(2.5),
                 Inches(2.3), Inches(1.0),
                 "기준년도 총 배출량",
                 f"{total_emissions/1e6:.1f}M tCO2e" if total_emissions else "N/A",
                 bg_color=WHITE, value_color=BLUE)

    # KPI: Net Zero 2050 NPV
    nz_npv = 0
    for d in comparison.get("npv_comparison", []):
        if d["scenario"] == "net_zero_2050":
            nz_npv = d["total_npv"]
            break
    _add_kpi_box(slide, LEFT_MARGIN + Inches(2.9), Inches(2.5),
                 Inches(2.3), Inches(1.0),
                 "Net Zero 2050 NPV",
                 f"-{_fmt_b(nz_npv)}",
                 bg_color=WHITE, value_color=RED)

    _add_textbox(slide, LEFT_MARGIN + Inches(0.3), Inches(3.8), Inches(5.0), Inches(0.7),
                 "탄소 비용 + 에너지 비용 + 매출 영향의 총합\n4개 NGFS 시나리오별 전환 비용 산출",
                 font_size=13, color=DARK_GRAY)

    # ── Right: Physical Risk Summary ──
    _add_rounded_rect(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(4.0),
                      LIGHT_ORANGE, shadow=True)
    _add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.06), ORANGE)
    _add_textbox(slide, Inches(7.3), Inches(1.8), Inches(5.0), Inches(0.5),
                 "물리적 리스크 요약", font_size=20, color=NAVY, bold=True)

    # KPI: Total EAL
    total_eal = sum(f["total_expected_annual_loss"]
                    for f in physical_data.get("facilities", []))
    _add_kpi_box(slide, Inches(7.3), Inches(2.5),
                 Inches(2.3), Inches(1.0),
                 "총 연간예상손실(EAL)",
                 f"${total_eal/1e6:.1f}M",
                 bg_color=WHITE, value_color=ORANGE)

    # KPI: High risk facility count
    summary = physical_data.get("overall_risk_summary", {})
    high_count = summary.get("High", 0)
    _add_kpi_box(slide, Inches(9.6), Inches(2.5),
                 Inches(2.3), Inches(1.0),
                 "고위험 시설 수",
                 f"{high_count}개",
                 bg_color=WHITE, value_color=RED)

    _add_textbox(slide, Inches(7.3), Inches(3.8), Inches(5.0), Inches(0.7),
                 "5대 기후재해 (홍수/태풍/폭염/가뭄/해수면)\nAnalytical v1 모델 + Open-Meteo 실측 데이터",
                 font_size=13, color=DARK_GRAY)

    # ── Bottom callout ──
    _add_rounded_rect(slide, LEFT_MARGIN, Inches(6.0), CONTENT_WIDTH, Inches(0.8),
                      NAVY, shadow=True)
    _add_textbox(slide, LEFT_MARGIN, Inches(6.1), CONTENT_WIDTH, Inches(0.5),
                 "핵심 질문: 기후 리스크가 존재하는가?  →  다음 단계 심층 분석 필요",
                 font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    _add_footer(slide, 5)


def slide_06_scenarios(prs):
    """Slide 6: NGFS 4 scenario analysis with rounded scenario cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "NGFS 4개 시나리오 분석", "시나리오별 탄소가격 경로 및 감축 목표")

    headers = ["시나리오", "2025", "2030", "2050", "감축 목표"]
    rows = []
    for sid in ["net_zero_2050", "below_2c", "delayed_transition", "current_policies"]:
        sc = SCENARIOS[sid]
        color = SCENARIO_COLORS[sid]
        rows.append([
            (sc["name"], color, True),
            f"${sc['carbon_price_2025']:.0f}",
            f"${sc['carbon_price_2030']:.0f}",
            f"${sc['carbon_price_2050']:.0f}",
            f"{sc['emissions_reduction_target']*100:.0f}%",
        ])

    _add_table(slide, LEFT_MARGIN, Inches(1.6), CONTENT_WIDTH, Inches(2.5),
               headers, rows,
               col_widths=[Inches(3.5), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.2)])

    # Scenario description cards (rounded + shadow)
    descs = [
        ("Net Zero 2050", "1.5C 목표, 즉각적이고 원활한 전환", RED),
        ("Below 2C", "2C 미만 목표, 점진적 전환", ORANGE),
        ("Delayed Transition", "2030년까지 정책 지연 후 급격한 전환", YELLOW),
        ("Current Policies", "현재 정책 유지, 제한적 추가 조치", GREEN),
    ]
    for i, (name, desc, color) in enumerate(descs):
        x = LEFT_MARGIN + (i % 2) * Inches(6.1)
        y = Inches(4.5) + (i // 2) * Inches(1.3)

        # Rounded card
        _add_rounded_rect(slide, x, y, Inches(5.7), Inches(1.1), LIGHT_GRAY, shadow=True)
        # Color accent left bar
        _add_shape(slide, x, y + Inches(0.1), Inches(0.06), Inches(0.9), color)

        _add_textbox(slide, x + Inches(0.25), y + Inches(0.1), Inches(5.2), Inches(0.4),
                     name, font_size=16, color=NAVY, bold=True)
        _add_textbox(slide, x + Inches(0.25), y + Inches(0.55), Inches(5.2), Inches(0.4),
                     desc, font_size=13, color=DARK_GRAY)

    _add_footer(slide, 6)


def slide_07_transition_risk(prs, comparison):
    """Slide 6: Transition risk - financial impact with chevron flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "전환 리스크 - 재무 영향 분석",
                   "배출경로 > ETS비용 > 에너지비용 > 매출영향 > NPV")

    # Chevron process flow
    flow_steps = ["배출 경로\n시뮬레이션", "ETS 비용\n산출", "에너지 비용\n증가분", "매출 영향\n분석", "NPV\n산출"]
    flow_colors = [BLUE, RGBColor(0x2B, 0x72, 0xE6), RGBColor(0x1B, 0x62, 0xD6),
                   RGBColor(0x0B, 0x52, 0xC6), NAVY]
    _add_chevron_flow(slide, flow_steps, Inches(1.55), flow_colors)

    # NPV comparison cards (rounded + shadow)
    npv_data = comparison["npv_comparison"]
    _add_textbox(slide, LEFT_MARGIN, Inches(2.8), Inches(11), Inches(0.5),
                 "시나리오별 총 NPV 영향", font_size=18, color=NAVY, bold=True)

    max_npv = max(abs(d["total_npv"]) for d in npv_data)
    for i, d in enumerate(npv_data):
        x = LEFT_MARGIN + i * Inches(3.05)
        y = Inches(3.35)
        sid = d["scenario"]
        color = SCENARIO_COLORS.get(sid, BLUE)
        bg = SCENARIO_BG_COLORS.get(sid, LIGHT_BLUE)

        # Rounded card + shadow
        _add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.2), bg, shadow=True)
        # Top accent
        _add_shape(slide, x + Inches(0.1), y, Inches(2.6), Inches(0.06), color)

        _add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.4), Inches(0.4),
                     d["scenario_name"], font_size=14, color=NAVY, bold=True,
                     alignment=PP_ALIGN.CENTER)

        _add_textbox(slide, x + Inches(0.2), y + Inches(0.65), Inches(2.4), Inches(0.6),
                     f"-{_fmt_b(d['total_npv'])}", font_size=32, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)

        _add_textbox(slide, x + Inches(0.2), y + Inches(1.35), Inches(2.4), Inches(0.4),
                     f"리스크 수준: {d['avg_risk_level']}", font_size=12, color=DARK_GRAY,
                     alignment=PP_ALIGN.CENTER)

        # Progress bar
        bar_pct = abs(d["total_npv"]) / max_npv if max_npv else 0
        bar_w = Emu(int(Inches(2.4) * bar_pct))
        _add_rounded_rect(slide, x + Inches(0.2), y + Inches(1.85),
                          bar_w, Inches(0.15), color, shadow=False)

    # Bottom note
    _add_textbox(slide, LEFT_MARGIN, Inches(5.85), CONTENT_WIDTH, Inches(0.5),
                 "* 시나리오에 따라 최대 2배 차이나는 전환 비용  |  할인율: 8% WACC  |  분석 기간: 2025~2050",
                 font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    _add_footer(slide, 7)


def slide_08_facility_ranking(prs, analysis_nz):
    """Slide 8: Facility risk ranking with fixed EMU values."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "시설별 리스크 순위", "Net Zero 2050 시나리오 기준 Top 5 고위험 시설")

    facs = sorted(analysis_nz["facilities"], key=lambda x: x["delta_npv"])
    risk_colors = {"High": RED, "Medium": ORANGE, "Low": GREEN}

    # Top 5 table
    headers = ["순위", "시설명", "섹터", "Delta NPV", "자산 대비 %", "리스크 등급"]
    rows = []
    for i, f in enumerate(facs[:5]):
        level = f["risk_level"]
        rows.append([
            str(i + 1),
            f["facility_name"],
            SECTOR_NAMES_KR.get(f["sector"], f["sector"]),
            f"-{_fmt_b(f['delta_npv'])}",
            f"{f['npv_as_pct_of_assets']:.1f}%",
            (level, risk_colors.get(level, BLACK), True),
        ])

    _add_table(slide, LEFT_MARGIN, Inches(1.6), CONTENT_WIDTH, Inches(2.8),
               headers, rows,
               col_widths=[Inches(1.0), Inches(3.0), Inches(2.0), Inches(2.0), Inches(1.8), Inches(1.9)])

    # Sector vulnerability ranking
    _add_textbox(slide, LEFT_MARGIN, Inches(4.6), Inches(11), Inches(0.5),
                 "섹터별 취약도 순위", font_size=18, color=NAVY, bold=True)

    sector_npv = {}
    for f in analysis_nz["facilities"]:
        s = f["sector"]
        sector_npv[s] = sector_npv.get(s, 0) + f["delta_npv"]

    sorted_sectors = sorted(sector_npv.items(), key=lambda x: x[1])

    # Horizontal bar chart
    max_val = abs(sorted_sectors[0][1]) if sorted_sectors else 1
    for i, (sector, npv) in enumerate(sorted_sectors):
        y = Inches(5.2) + i * Inches(0.27)
        bar_pct = abs(npv) / max_val
        bar_w = Emu(int(max(Inches(0.1), Inches(8.0) * bar_pct)))

        _add_textbox(slide, LEFT_MARGIN, y - Inches(0.02), Inches(2.0), Inches(0.3),
                     SECTOR_NAMES_KR.get(sector, sector), font_size=10, color=DARK_GRAY,
                     alignment=PP_ALIGN.RIGHT)
        _add_rounded_rect(slide, Inches(3.0), y, bar_w, Inches(0.2), BLUE, shadow=False)
        _add_textbox(slide, Emu(int(Inches(3.0) + bar_w + Inches(0.1))), y - Inches(0.02),
                     Inches(2.0), Inches(0.3),
                     f"-{_fmt_b(npv)}", font_size=10, color=DARK_GRAY)

    # Risk grade legend
    _add_rounded_rect(slide, LEFT_MARGIN, Inches(7.0), CONTENT_WIDTH, Inches(0.35),
                      LIGHT_GRAY, shadow=False)
    _add_textbox(slide, Inches(1.0), Inches(7.02), Inches(11.0), Inches(0.3),
                 "리스크 등급 기준:  High = NPV가 자산의 -15% 이상  |  Medium = -5%~-15%  |  Low = -5% 미만",
                 font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    _add_footer(slide, 8)


def slide_09_physical_risk(prs, physical_data):
    """Slide 9: Physical risk assessment with real data from physical_risk.py."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "물리적 리스크 평가", "기후재해 유형별 시설 취약도 평가")

    # 5 hazard cards (text-based, no emoji)
    hazards = [
        ("홍수", "집중호우/하천범람\n연안 침수 위험", BLUE),
        ("태풍", "강풍/폭풍해일\n시설 물리적 손상", RED),
        ("폭염", "냉각 효율 저하\n작업자 안전/생산성", ORANGE),
        ("가뭄", "용수 공급 차질\n냉각수 부족", YELLOW),
        ("해수면", "연안 시설 침수\n장기 입지 리스크", RGBColor(0x06, 0xB6, 0xD4)),
    ]

    for i, (title, desc, color) in enumerate(hazards):
        x = Inches(0.5) + i * Inches(2.5)
        y = Inches(1.7)
        _add_rounded_rect(slide, x, y, Inches(2.2), Inches(1.6), LIGHT_BLUE, shadow=True)
        # Top accent
        _add_shape(slide, x, y, Inches(2.2), Inches(0.05), color)
        _add_textbox(slide, x + Inches(0.15), y + Inches(0.15), Inches(1.9), Inches(0.4),
                     title, font_size=16, color=NAVY, bold=True,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.15), y + Inches(0.6), Inches(1.9), Inches(0.8),
                     desc, font_size=12, color=DARK_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # EAL Top 5 table from real data
    _add_textbox(slide, LEFT_MARGIN, Inches(3.5), Inches(11), Inches(0.5),
                 "시설별 연간 예상손실(EAL) Top 5", font_size=18, color=NAVY, bold=True)

    fac_results = sorted(physical_data["facilities"],
                         key=lambda x: x["total_expected_annual_loss"], reverse=True)

    headers = ["순위", "시설명", "위치", "위험 등급", "EAL ($M)", "주요 재해"]
    rows = []
    for i, fac in enumerate(fac_results[:5]):
        # Find the highest risk hazard
        top_hazard = max(fac["hazards"], key=lambda h: h["potential_loss"])
        rows.append([
            str(i + 1),
            fac["facility_name"],
            fac["location"],
            (fac["overall_risk_level"],
             {"High": RED, "Medium": ORANGE, "Low": GREEN}.get(fac["overall_risk_level"], BLACK),
             True),
            f"${fac['total_expected_annual_loss']/1e6:.1f}",
            HAZARD_NAMES_KR.get(top_hazard["hazard_type"], top_hazard["hazard_type"]),
        ])

    _add_table(slide, LEFT_MARGIN, Inches(4.0), CONTENT_WIDTH, Inches(2.0),
               headers, rows,
               col_widths=[Inches(0.8), Inches(2.5), Inches(2.0), Inches(1.8), Inches(2.2), Inches(2.4)])

    # Risk distribution summary
    summary = physical_data["overall_risk_summary"]
    _add_rounded_rect(slide, LEFT_MARGIN, Inches(6.2), CONTENT_WIDTH, Inches(0.8), NAVY, shadow=True)

    _add_rich_text(slide, Inches(1.0), Inches(6.3), Inches(5.0), Inches(0.5),
                   [
                       ("시설 리스크 분포:  ", 14, WHITE, True),
                       (f"High {summary.get('High', 0)}개", 14, LIGHT_RED, True),
                       ("  |  ", 14, MID_GRAY, False),
                       (f"Medium {summary.get('Medium', 0)}개", 14, LIGHT_ORANGE, True),
                       ("  |  ", 14, MID_GRAY, False),
                       (f"Low {summary.get('Low', 0)}개", 14, LIGHT_GREEN, True),
                   ])

    _add_textbox(slide, Inches(7.0), Inches(6.35), Inches(5.5), Inches(0.4),
                 "Analytical v1 모델 | KMA 30년 통계 + IPCC AR6 + Open-Meteo API",
                 font_size=12, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

    _add_footer(slide, 9)


def slide_10_stage3_cashflow(prs):
    """Slide 10: Stage 3 cashflow impact vision (methodology only, not yet implemented)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Stage 3: 캐시플로우 영향 분석",
                   "기후 리스크가 기업 현금흐름에 미치는 영향 (개발 예정)")

    # Chevron process flow
    flow_steps = ["기존 FCF\n추정", "전환 리스크\n반영", "물리적 리스크\n반영", "할인율\n조정", "기후조정\nNPV"]
    flow_colors = [BLUE, RGBColor(0x2B, 0x72, 0xE6), ORANGE,
                   RGBColor(0xE0, 0x58, 0x10), RED]
    _add_chevron_flow(slide, flow_steps, Inches(1.6), flow_colors)

    # 3 output cards
    outputs = [
        ("시나리오별 조정 FCF", "탄소비용, EAL, 자산손상 반영\n시나리오별 현금흐름 변동 추정", BLUE, LIGHT_BLUE),
        ("Climate-adjusted WACC", "시나리오별 리스크 프리미엄 할인율\n기후 리스크 반영 자본비용", ORANGE, LIGHT_ORANGE),
        ("Climate VaR", "최악 시나리오 최대 가치 하락폭\n포트폴리오 전체 리스크 규모", RED, LIGHT_RED),
    ]

    for i, (title, desc, color, bg_color) in enumerate(outputs):
        x = LEFT_MARGIN + i * Inches(4.1)
        y = Inches(2.9)

        _add_rounded_rect(slide, x, y, Inches(3.7), Inches(2.2), bg_color, shadow=True)
        _add_shape(slide, x, y, Inches(3.7), Inches(0.06), color)

        _add_textbox(slide, x + Inches(0.25), y + Inches(0.25), Inches(3.2), Inches(0.5),
                     title, font_size=18, color=NAVY, bold=True)
        _add_textbox(slide, x + Inches(0.25), y + Inches(0.85), Inches(3.2), Inches(1.0),
                     desc, font_size=14, color=DARK_GRAY)

    # "개발 예정" badge
    badge = _add_rounded_rect(slide, Inches(5.5), Inches(5.4), Inches(2.3), Inches(0.4),
                               MID_GRAY, shadow=False)
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "2026 H2 개발 예정"
    _set_font(run, size=13, color=WHITE, bold=True)

    # Bottom callout
    _add_rounded_rect(slide, LEFT_MARGIN, Inches(6.0), CONTENT_WIDTH, Inches(0.8),
                      NAVY, shadow=True)
    _add_textbox(slide, LEFT_MARGIN, Inches(6.1), CONTENT_WIDTH, Inches(0.5),
                 "Stage 1 + 2의 분석 결과를 통합한 재무 모델",
                 font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    _add_footer(slide, 10)


def slide_11_esg(prs):
    """Slide 11: ESG compliance one-stop support - text-based status."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "ESG 공시 원스톱 지원", "TCFD / ISSB / KSSB 3대 프레임워크 자동 분석")

    # Get actual scores
    frameworks = {}
    for fid in ["tcfd", "issb", "kssb"]:
        frameworks[fid] = assess_framework(fid)

    # Framework cards (rounded + shadow)
    cards = [
        ("TCFD", "tcfd", BLUE),
        ("ISSB (IFRS S2)", "issb", ORANGE),
        ("KSSB", "kssb", GREEN),
    ]
    for i, (name, fid, color) in enumerate(cards):
        x = LEFT_MARGIN + i * Inches(4.1)
        y = Inches(1.7)
        fw = frameworks[fid]

        _add_rounded_rect(slide, x, y, Inches(3.7), Inches(2.5), LIGHT_GRAY, shadow=True)
        _add_shape(slide, x, y, Inches(3.7), Inches(0.06), color)

        _add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.3), Inches(0.4),
                     name, font_size=18, color=NAVY, bold=True)

        # Score
        score = fw["overall_score"]
        level = fw["compliance_level"]
        _add_textbox(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.3), Inches(0.6),
                     f"{score:.0f}점 / 100점", font_size=24, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.2), y + Inches(1.3), Inches(3.3), Inches(0.3),
                     f"준수 수준: {level}", font_size=14, color=DARK_GRAY,
                     alignment=PP_ALIGN.CENTER)

        # Text-based checklist summary (no emoji)
        checklist = fw["checklist"]
        compliant = sum(1 for c in checklist if c["status"] == "compliant")
        partial = sum(1 for c in checklist if c["status"] == "partial")
        non_compliant = sum(1 for c in checklist if c["status"] == "non_compliant")

        _add_rich_text(slide, x + Inches(0.15), y + Inches(1.75), Inches(3.4), Inches(0.6),
                       [
                           ("준수 ", 11, DARK_GRAY, False),
                           (f"{compliant}", 11, GREEN, True),
                           ("  |  부분 ", 11, DARK_GRAY, False),
                           (f"{partial}", 11, ORANGE, True),
                           ("  |  미준수 ", 11, DARK_GRAY, False),
                           (f"{non_compliant}", 11, RED, True),
                       ],
                       alignment=PP_ALIGN.CENTER)

    # Key features (text-based, bullet list)
    features = [
        "자동 준수도 점검 (체크리스트 + 점수)",
        "공시 서술문 초안 자동 생성 (거버넌스/전략/리스크관리/지표)",
        "개선 권고사항 도출 및 우선순위 제시",
        "분석 결과 > 공시 데이터 자동 연계",
    ]
    _add_textbox(slide, LEFT_MARGIN, Inches(4.5), Inches(11), Inches(0.5),
                 "주요 기능", font_size=18, color=NAVY, bold=True)
    _add_bullet_list(slide, LEFT_MARGIN, Inches(5.1), CONTENT_WIDTH, Inches(2.0),
                     features, font_size=14, bullet_char="- ")

    _add_footer(slide, 11)


def _deleted_slide_10_dashboard(prs):
    """Slide 10: Dashboard preview with rounded cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "대시보드 미리보기", "웹 기반 인터랙티브 분석 대시보드")

    # 6 dashboard pages
    pages = [
        ("Overview", "종합 현황 대시보드\n핵심 KPI 한눈에 확인"),
        ("Emissions", "배출경로 시뮬레이션\n시나리오별 감축 경로 비교"),
        ("Financial Impact", "NPV 비교 차트\n비용 워터폴 분석"),
        ("Facility Risk", "시설별 리스크 히트맵\n순위 테이블"),
        ("Physical Risk", "지역별 재해 위험도\n시설 매핑"),
        ("ESG Compliance", "프레임워크별 레이더 차트\n체크리스트/권고사항"),
    ]

    for i, (title, desc) in enumerate(pages):
        x = LEFT_MARGIN + (i % 3) * Inches(4.1)
        y = Inches(1.7) + (i // 3) * Inches(2.3)

        # Dashboard page mock (rounded)
        _add_rounded_rect(slide, x, y, Inches(3.7), Inches(1.9), LIGHT_GRAY, shadow=True)
        # Top bar
        _add_shape(slide, x, y, Inches(3.7), Inches(0.4), NAVY)
        _add_textbox(slide, x + Inches(0.15), y + Inches(0.05), Inches(3.4), Inches(0.3),
                     title, font_size=12, color=WHITE, bold=True)
        # Description
        _add_textbox(slide, x + Inches(0.2), y + Inches(0.55), Inches(3.3), Inches(1.2),
                     desc, font_size=13, color=DARK_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # Bottom features (rounded)
    _add_rounded_rect(slide, LEFT_MARGIN, Inches(6.3), CONTENT_WIDTH, Inches(0.9),
                      LIGHT_BLUE, shadow=True)
    feature_items = [
        ("시나리오 전환", "드롭다운 선택,\n실시간 결과 변경"),
        ("인터랙티브 차트", "Recharts 기반\n줌/필터/호버 지원"),
        ("웹 기반 접속", "별도 설치 없이\n브라우저에서 즉시 사용"),
    ]
    for i, (ft, fd) in enumerate(feature_items):
        x = Inches(1.2) + i * Inches(4.0)
        _add_textbox(slide, x, Inches(6.35), Inches(3.5), Inches(0.35),
                     ft, font_size=14, color=NAVY, bold=True)
        _add_textbox(slide, x, Inches(6.7), Inches(3.5), Inches(0.4),
                     fd, font_size=11, color=DARK_GRAY)

    _add_footer(slide, 10)


def slide_12_roadmap(prs):
    """Slide 12: Stage-based roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "단계 기반 로드맵", "분석 심화 단계별 기능 확장 계획")

    phases = [
        ("Stage 1 + 2", "현재 제공", BLUE, [
            "전환 리스크 시나리오 분석 (NGFS 4개)",
            "K-ETS 무상할당 분석",
            "물리적 리스크 평가 (Analytical v1)",
            "Open-Meteo 실측 데이터 연동",
            "ESG 공시 (TCFD/ISSB/KSSB)",
            "인터랙티브 웹 대시보드",
        ]),
        ("Stage 3", "2026 H2", ORANGE, [
            "현금흐름 영향 모델 (DCF)",
            "전환+물리적 통합 재무분석",
            "Climate VaR 산출",
            "시나리오별 WACC 조정",
            "MACC 감축 비용 곡선",
        ]),
        ("확장 기능", "2027", GREEN, [
            "고객 데이터 직접 업로드",
            "DB 연동 (Oracle, PostgreSQL)",
            "API 연동 (ERP, 탄소회계)",
            "산업별 벤치마크 비교",
            "PDF 리포트 자동 생성",
        ]),
    ]

    for i, (phase, timing, color, items) in enumerate(phases):
        x = LEFT_MARGIN + i * Inches(4.1)
        y = Inches(1.6)

        # Phase header (rounded top)
        _add_rounded_rect(slide, x, y, Inches(3.7), Inches(0.7), color, shadow=False)
        _add_textbox(slide, x + Inches(0.2), y + Inches(0.05), Inches(2.5), Inches(0.35),
                     phase, font_size=20, color=WHITE, bold=True)
        _add_textbox(slide, x + Inches(0.2), y + Inches(0.35), Inches(2.5), Inches(0.3),
                     timing, font_size=13, color=WHITE)

        # Items (rounded card)
        _add_rounded_rect(slide, x, y + Inches(0.7), Inches(3.7), Inches(4.0),
                          LIGHT_GRAY, shadow=True)
        _add_bullet_list(slide, x + Inches(0.15), y + Inches(0.9), Inches(3.4), Inches(3.6),
                         items, font_size=13, bullet_char="- ")

    # Timeline arrow
    _add_shape(slide, LEFT_MARGIN, Inches(6.5), CONTENT_WIDTH, Inches(0.06), BLUE)
    for i in range(3):
        x = Inches(2.8) + i * Inches(4.1)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(6.38), Inches(0.25), Inches(0.25))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background()

    _add_footer(slide, 12)


def slide_13_contact(prs):
    """Slide 13: Contact / deployment guide with gradient background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, '0F172A', '1E293B')

    # Title
    _add_textbox(slide, Inches(1.0), Inches(0.8), Inches(11), Inches(0.8),
                 "도입 안내", font_size=36, color=WHITE, bold=True)
    _add_shape(slide, Inches(1.0), Inches(1.5), Inches(2.0), Inches(0.04), BLUE)

    # Deployment options (text-based, no emoji)
    options = [
        ("클라우드 (SaaS)", "별도 설치 없이 즉시 사용\n자동 업데이트, 유지보수 포함\n월 구독 요금제", BLUE),
        ("온프레미스", "고객사 서버에 직접 설치\n데이터 외부 유출 차단\n커스텀 보안 정책 적용", ORANGE),
        ("커스터마이징", "고객사 데이터 맞춤 설정\n산업별 파라미터 조정\n기존 시스템 연동 지원", GREEN),
    ]

    for i, (title, desc, accent) in enumerate(options):
        x = Inches(1.0) + i * Inches(4.1)
        y = Inches(2.2)

        _add_rounded_rect(slide, x, y, Inches(3.7), Inches(2.4), NAVY, shadow=True)
        # Top accent
        _add_shape(slide, x, y, Inches(3.7), Inches(0.05), accent)
        _add_textbox(slide, x + Inches(0.25), y + Inches(0.2), Inches(3.2), Inches(0.5),
                     title, font_size=18, color=accent, bold=True)
        _add_textbox(slide, x + Inches(0.25), y + Inches(0.8), Inches(3.2), Inches(1.4),
                     desc, font_size=14, color=MID_GRAY)

    # Demo CTA
    _add_rounded_rect(slide, Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.0), BLUE, shadow=True)
    _add_textbox(slide, Inches(1.0), Inches(5.1), Inches(11.3), Inches(0.5),
                 "즉시 체험 가능: 데모 사이트 또는 로컬 설치 (Docker 기반 원클릭 배포)",
                 font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.0), Inches(5.55), Inches(11.3), Inches(0.4),
                 "npm run dev (프론트엔드) + uvicorn (백엔드) > 5분 내 로컬 환경 구축",
                 font_size=13, color=DARK_NAVY, alignment=PP_ALIGN.CENTER)

    # Contact
    _add_textbox(slide, Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.5),
                 "문의: climate-risk@example.com  |  Tel: 02-XXX-XXXX",
                 font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1.0), Inches(6.9), Inches(11.3), Inches(0.4),
                 "Climate Risk Platform  |  Confidential",
                 font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════

def main():
    print("=" * 60)
    print("  Climate Risk Platform - PPT 생성 (3단계 분석 흐름)")
    print("=" * 60)

    # Load real data from backend
    print("\n[1/5] 백엔드 데이터 로딩...")
    facilities = get_all_facilities()
    print(f"  > 시설 {len(facilities)}개 로드 완료")

    print("[2/5] 시나리오 분석 실행...")
    analysis_nz = analyse_scenario("net_zero_2050")
    comparison = compare_scenarios()
    print(f"  > Net Zero NPV: ${analysis_nz['total_npv']:,.0f}")
    print(f"  > 4개 시나리오 비교 완료")

    print("[3/5] 물리적 리스크 분석 실행...")
    physical_data = assess_physical_risk()
    print(f"  > {physical_data['total_facilities']}개 시설 물리적 리스크 평가 완료")
    print(f"  > 리스크 분포: {physical_data['overall_risk_summary']}")

    print("[4/5] PPT 슬라이드 생성 중 (13장)...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_01_cover(prs)
    print("  > Slide  1: 표지")
    slide_02_problem(prs)
    print("  > Slide  2: 기업이 직면한 문제")
    slide_03_service_stages(prs)
    print("  > Slide  3: 3단계 분석 방법론")
    slide_04_portfolio(prs, facilities)
    print("  > Slide  4: 분석 포트폴리오")
    slide_05_stage1_overview(prs, comparison, physical_data)
    print("  > Slide  5: Stage 1 기초 분석 요약")
    slide_06_scenarios(prs)
    print("  > Slide  6: NGFS 4개 시나리오")
    slide_07_transition_risk(prs, comparison)
    print("  > Slide  7: 전환 리스크 재무영향")
    slide_08_facility_ranking(prs, analysis_nz)
    print("  > Slide  8: 시설별 리스크 순위")
    slide_09_physical_risk(prs, physical_data)
    print("  > Slide  9: 물리적 리스크 평가")
    slide_10_stage3_cashflow(prs)
    print("  > Slide 10: Stage 3 캐시플로우 비전")
    slide_11_esg(prs)
    print("  > Slide 11: ESG 공시 지원")
    slide_12_roadmap(prs)
    print("  > Slide 12: 단계 기반 로드맵")
    slide_13_contact(prs)
    print("  > Slide 13: 도입 안내")

    print("[5/5] 파일 저장 중...")
    output_dir = os.path.join(os.path.dirname(__file__), "outputs")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "climate_risk_platform_overview.pptx")
    prs.save(output_path)
    print(f"\nPPT 생성 완료: {output_path}")
    print(f"  슬라이드 수: {len(prs.slides)}")
    print(f"  폰트: {LATIN_FONT} (Latin) + {EA_FONT} (EA)")
    print(f"  플랫폼: {platform.system()}")
    print("=" * 60)


if __name__ == "__main__":
    main()
