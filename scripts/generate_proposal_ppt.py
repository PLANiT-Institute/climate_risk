#!/usr/bin/env python3
"""Generate a 20-slide client proposal PPT for climate risk disclosure advisory.

Slides:
  Part 1 (1-4):  Why Climate Disclosure Matters
  Part 2 (5-12): Your Climate Risk Profile
  Part 3 (13-16): ESG Disclosure Readiness
  Part 4 (17-20): Our Solution

Usage:
    source backend/venv/bin/activate
    python scripts/generate_proposal_ppt.py
    python scripts/generate_proposal_ppt.py --scenario below_2c --pricing global
    python scripts/generate_proposal_ppt.py --output /path/to/output.pptx
"""

import argparse
import math
import os
import platform
import sys
import tempfile
from collections import Counter
from datetime import date, datetime

# Add backend to path so we can import services
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np

from app.data.sample_facilities import get_all_facilities
from app.services.transition_risk import (
    analyse_scenario,
    compare_scenarios,
    get_summary,
)
from app.services.physical_risk import assess_physical_risk
from app.services.esg_compliance import assess_framework, get_disclosure_data
from app.core.config import (
    SCENARIOS,
    NGFS_CARBON_PRICE_PATHS,
    KETS_ALLOCATION_RATIOS,
    REGULATORY_DEADLINES,
)

# ══════════════════════════════════════════════════════════════════════
# COLOR PALETTE
# ══════════════════════════════════════════════════════════════════════

NAVY = RGBColor(0x1E, 0x29, 0x3B)
DARK_NAVY = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
ALT_ROW_GRAY = RGBColor(0xF8, 0xFA, 0xFC)
MID_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY = RGBColor(0x47, 0x55, 0x69)
BLACK = RGBColor(0x1E, 0x1E, 0x1E)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
YELLOW = RGBColor(0xEA, 0xB3, 0x08)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
LIGHT_GREEN = RGBColor(0xDC, 0xFC, 0xE7)
LIGHT_RED = RGBColor(0xFE, 0xE2, 0xE2)
LIGHT_ORANGE = RGBColor(0xFF, 0xED, 0xD5)
LIGHT_YELLOW = RGBColor(0xFE, 0xF9, 0xC3)
TEAL = RGBColor(0x06, 0xB6, 0xD4)

SCENARIO_COLORS = {
    "net_zero_2050": RED,
    "below_2c": ORANGE,
    "delayed_transition": YELLOW,
    "current_policies": GREEN,
}
SCENARIO_BG_COLORS = {
    "net_zero_2050": LIGHT_RED,
    "below_2c": LIGHT_ORANGE,
    "delayed_transition": LIGHT_YELLOW,
    "current_policies": LIGHT_GREEN,
}

# Matplotlib hex equivalents
MPL_NAVY = "#1E293B"
MPL_BLUE = "#3B82F6"
MPL_RED = "#EF4444"
MPL_ORANGE = "#F97316"
MPL_YELLOW = "#EAB308"
MPL_GREEN = "#22C55E"
MPL_TEAL = "#06B6D4"
MPL_GRAY = "#94A3B8"
MPL_DARK = "#475569"
MPL_LIGHT_BG = "#F8FAFC"
MPL_SCENARIO_COLORS = {
    "net_zero_2050": MPL_RED,
    "below_2c": MPL_ORANGE,
    "delayed_transition": MPL_YELLOW,
    "current_policies": MPL_GREEN,
}

# ── Slide dimensions (16:9) ─────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Layout constants ─────────────────────────────────────────────────
LEFT_MARGIN = Inches(0.8)
RIGHT_MARGIN = Inches(0.8)
CONTENT_WIDTH = Inches(11.7)

# ── Font strategy ────────────────────────────────────────────────────
_IS_MAC = platform.system() == "Darwin"
LATIN_FONT = "Calibri"
EA_FONT = "Apple SD Gothic Neo" if _IS_MAC else "Malgun Gothic"

# ── Name maps ────────────────────────────────────────────────────────
SECTOR_NAMES_KR = {
    "steel": "철강",
    "petrochemical": "석유화학",
    "automotive": "자동차",
    "electronics": "반도체/전자",
    "utilities": "발전",
    "cement": "시멘트",
    "shipping": "해운",
    "oil_gas": "정유",
}
HAZARD_NAMES_KR = {
    "flood": "홍수",
    "typhoon": "태풍",
    "heatwave": "폭염",
    "drought": "가뭄",
    "sea_level_rise": "해수면",
}
SCENARIO_NAMES_KR = {
    "net_zero_2050": "넷제로 2050",
    "below_2c": "2°C 이하",
    "delayed_transition": "지연 전환",
    "current_policies": "현행 정책",
}


# ══════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS (reused from generate_ppt.py + new)
# ══════════════════════════════════════════════════════════════════════

def _set_font(run, name=LATIN_FONT, ea_name=EA_FONT, size=18, bold=False, color=BLACK):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for ea in rPr.findall(qn("a:ea")):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", ea_name)


def _add_shadow(shape):
    spPr = shape._element.spPr
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    shadow = etree.SubElement(effectLst, qn("a:outerShdw"))
    shadow.set("blurRad", "50800")
    shadow.set("dist", "25400")
    shadow.set("dir", "5400000")
    shadow.set("algn", "bl")
    shadow.set("rotWithShape", "0")
    srgb = etree.SubElement(shadow, qn("a:srgbClr"))
    srgb.set("val", "000000")
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", "20000")


def _set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_gradient_bg(slide, color1_hex, color2_hex, angle=5400000):
    cSld = slide._element
    bg_elem = cSld.find(qn("p:bg"))
    if bg_elem is None:
        bg_elem = etree.SubElement(cSld, qn("p:bg"))
        cSld.insert(0, bg_elem)
    for old in bg_elem.findall(qn("p:bgPr")):
        bg_elem.remove(old)
    bgPr = etree.SubElement(bg_elem, qn("p:bgPr"))
    gradFill = etree.SubElement(bgPr, qn("a:gradFill"))
    gradFill.set("flip", "none")
    gradFill.set("rotWithShape", "1")
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    gs1 = etree.SubElement(gsLst, qn("a:gs"))
    gs1.set("pos", "0")
    srgb1 = etree.SubElement(gs1, qn("a:srgbClr"))
    srgb1.set("val", color1_hex)
    gs2 = etree.SubElement(gsLst, qn("a:gs"))
    gs2.set("pos", "100000")
    srgb2 = etree.SubElement(gs2, qn("a:srgbClr"))
    srgb2.set("val", color2_hex)
    lin = etree.SubElement(gradFill, qn("a:lin"))
    lin.set("ang", str(angle))
    lin.set("scaled", "1")
    etree.SubElement(bgPr, qn("a:effectLst"))


def _add_rounded_rect(slide, left, top, width, height, fill_color, shadow=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if shadow:
        _add_shadow(shape)
    return shape


def _add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, size=font_size, color=color, bold=bold)
    return txBox


def _add_rich_text(slide, left, top, width, height, segments, alignment=PP_ALIGN.LEFT):
    """segments = [(text, font_size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    for text, font_size, color, bold in segments:
        run = p.add_run()
        run.text = text
        _set_font(run, size=font_size, color=color, bold=bold)
    return txBox


def _add_title_bar(slide, title_text, subtitle_text=None):
    _add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), NAVY)
    _add_textbox(slide, LEFT_MARGIN, Inches(0.15), Inches(10), Inches(0.6),
                 title_text, font_size=28, color=WHITE, bold=True)
    if subtitle_text:
        _add_textbox(slide, LEFT_MARGIN, Inches(0.7), Inches(10), Inches(0.4),
                     subtitle_text, font_size=14, color=MID_GRAY)
    _add_shape(slide, Inches(0), Inches(1.2), SLIDE_WIDTH, Inches(0.04), BLUE)


def _add_kpi_box(slide, left, top, width, height, label, value,
                 bg_color=LIGHT_BLUE, value_color=BLUE):
    _add_rounded_rect(slide, left, top, width, height, bg_color, shadow=True)
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.4),
                 label, font_size=11, color=DARK_GRAY)
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.45), width - Inches(0.3), Inches(0.5),
                 value, font_size=28, color=value_color, bold=True)


def _add_table(slide, left, top, width, height, headers, rows, col_widths=None,
               header_bg=NAVY, header_fg=WHITE, stripe_color=ALT_ROW_GRAY):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        _style_cell(cell, font_size=12, bold=True, color=header_fg, bg_color=header_bg)
    for i, row in enumerate(rows):
        bg = stripe_color if i % 2 == 1 else WHITE
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            if isinstance(val, tuple):
                cell.text = val[0]
                _style_cell(cell, font_size=11, color=val[1],
                            bg_color=bg, bold=val[2] if len(val) > 2 else False)
            else:
                cell.text = str(val)
                _style_cell(cell, font_size=11, color=BLACK, bg_color=bg)
    return table_shape


def _style_cell(cell, font_size=11, bold=False, color=BLACK, bg_color=WHITE):
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg_color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            _set_font(run, size=font_size, bold=bold, color=color)
        if not p.runs:
            run = p.add_run()
            run.text = p.text
            _set_font(run, size=font_size, bold=bold, color=color)
        p.alignment = PP_ALIGN.CENTER


def _add_footer(slide, page_num):
    _add_shape(slide, LEFT_MARGIN, Inches(7.15), CONTENT_WIDTH, Inches(0.02), MID_GRAY)
    _add_textbox(slide, Inches(12.0), Inches(7.18), Inches(1.0), Inches(0.3),
                 str(page_num), font_size=10, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)
    _add_textbox(slide, LEFT_MARGIN, Inches(7.18), Inches(6.0), Inches(0.3),
                 "Climate Risk Advisory  |  Confidential", font_size=10, color=MID_GRAY)


def _fmt_b(val):
    return f"${abs(val)/1e9:.1f}B"


def _fmt_m(val):
    return f"{val/1e6:.1f}M"


def _fmt_pct(val):
    return f"{val:.0f}%"


# ── New helpers ──────────────────────────────────────────────────────

def _add_section_divider(prs, section_number, title_text, subtitle_text=""):
    """Full-slide gradient section break."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, "0F172A", "1E293B", angle=2700000)
    # Large section number
    _add_textbox(slide, Inches(1.2), Inches(1.5), Inches(3.0), Inches(2.0),
                 f"{section_number:02d}", font_size=96, color=BLUE, bold=True)
    # Accent line
    _add_shape(slide, Inches(1.2), Inches(3.5), Inches(2.0), Inches(0.06), BLUE)
    # Title
    _add_textbox(slide, Inches(1.2), Inches(3.8), Inches(10), Inches(1.0),
                 title_text, font_size=40, color=WHITE, bold=True)
    # Subtitle
    if subtitle_text:
        _add_textbox(slide, Inches(1.2), Inches(5.0), Inches(10), Inches(0.6),
                     subtitle_text, font_size=18, color=MID_GRAY)
    return slide


def _add_bullet_list(slide, left, top, width, height, items, font_size=14,
                     color=DARK_GRAY, bullet_color=BLUE):
    """Add a bulleted text list to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        # bullet marker
        run_bullet = p.add_run()
        run_bullet.text = "\u25CF  "
        _set_font(run_bullet, size=font_size, color=bullet_color, bold=False)
        # text
        run_text = p.add_run()
        run_text.text = item
        _set_font(run_text, size=font_size, color=color, bold=False)
    return txBox


def _embed_chart(slide, fig, left, top, width, height):
    """Save matplotlib figure to temp file and embed in slide."""
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    try:
        fig.savefig(tmp.name, dpi=200, bbox_inches="tight",
                    facecolor=fig.get_facecolor(), edgecolor="none")
        plt.close(fig)
        slide.shapes.add_picture(tmp.name, left, top, width, height)
    finally:
        try:
            os.unlink(tmp.name)
        except OSError:
            pass


def _add_card(slide, left, top, width, height, title, body_lines,
              accent_color=BLUE, bg_color=WHITE):
    """Add a card with accent top line, title, and body lines."""
    _add_rounded_rect(slide, left, top, width, height, bg_color, shadow=True)
    _add_shape(slide, left + Inches(0.05), top, width - Inches(0.1), Inches(0.05), accent_color)
    _add_textbox(slide, left + Inches(0.2), top + Inches(0.15),
                 width - Inches(0.4), Inches(0.4),
                 title, font_size=15, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    y_off = Inches(0.55)
    for line in body_lines:
        _add_textbox(slide, left + Inches(0.15), top + y_off,
                     width - Inches(0.3), Inches(0.3),
                     line, font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
        y_off += Inches(0.28)


# ══════════════════════════════════════════════════════════════════════
# CHART GENERATION FUNCTIONS (matplotlib → PNG)
# ══════════════════════════════════════════════════════════════════════

def _mpl_setup(figsize=(10, 5)):
    """Common matplotlib figure setup."""
    # Set Korean-capable font for matplotlib
    if _IS_MAC:
        plt.rcParams["font.family"] = "AppleGothic"
    else:
        plt.rcParams["font.family"] = "Malgun Gothic"
    plt.rcParams["axes.unicode_minus"] = False
    fig, ax = plt.subplots(figsize=figsize)
    fig.patch.set_facecolor(MPL_LIGHT_BG)
    ax.set_facecolor(MPL_LIGHT_BG)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color(MPL_GRAY)
    ax.spines["bottom"].set_color(MPL_GRAY)
    ax.tick_params(colors=MPL_DARK, labelsize=9)
    return fig, ax


def chart_scenario_npv_comparison(comparison_data):
    """Horizontal bar chart of total NPV impact by scenario."""
    fig, ax = _mpl_setup(figsize=(10, 4.5))
    npv_data = comparison_data["npv_comparison"]
    scenarios = [d["scenario_name"] for d in npv_data]
    values = [abs(d["total_npv"]) / 1e9 for d in npv_data]
    colors = [MPL_SCENARIO_COLORS.get(d["scenario"], MPL_BLUE) for d in npv_data]

    bars = ax.barh(scenarios, values, color=colors, height=0.6, edgecolor="white")
    ax.set_xlabel("Total NPV Impact ($B)", fontsize=10, color=MPL_DARK)
    ax.set_title("시나리오별 전환 리스크 NPV 영향", fontsize=13, fontweight="bold",
                 color=MPL_NAVY, pad=12)
    ax.invert_yaxis()
    for bar, v in zip(bars, values):
        ax.text(bar.get_width() + max(values) * 0.02, bar.get_y() + bar.get_height() / 2,
                f"${v:.1f}B", va="center", fontsize=10, fontweight="bold", color=MPL_DARK)
    ax.set_xlim(0, max(values) * 1.2)
    fig.tight_layout()
    return fig


def chart_cost_waterfall(summary_data):
    """Stacked waterfall chart showing cost breakdown."""
    fig, ax = _mpl_setup(figsize=(10, 4.5))
    breakdown = summary_data["cost_breakdown"]
    categories = ["탄소 비용", "에너지 비용", "매출 영향", "전환 OPEX", "총 NPV"]
    vals = [
        abs(breakdown["carbon_cost"]) / 1e9,
        abs(breakdown["energy_cost_increase"]) / 1e9,
        abs(breakdown["revenue_impact"]) / 1e9,
        abs(breakdown["transition_opex"]) / 1e9,
        abs(summary_data["total_npv"]) / 1e9,
    ]
    colors_list = [MPL_RED, MPL_ORANGE, MPL_YELLOW, MPL_TEAL, MPL_NAVY]
    bars = ax.bar(categories, vals, color=colors_list, width=0.6, edgecolor="white")
    ax.set_ylabel("Cost ($B)", fontsize=10, color=MPL_DARK)
    ax.set_title("전환 비용 구성 분석", fontsize=13, fontweight="bold",
                 color=MPL_NAVY, pad=12)
    for bar, v in zip(bars, vals):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + max(vals) * 0.02,
                f"${v:.1f}B", ha="center", fontsize=9, fontweight="bold", color=MPL_DARK)
    ax.set_ylim(0, max(vals) * 1.25)
    fig.tight_layout()
    return fig


def chart_carbon_price_paths():
    """Line chart showing 4-scenario carbon price trajectories."""
    fig, ax = _mpl_setup(figsize=(10, 4.5))
    for sid, path in NGFS_CARBON_PRICE_PATHS.items():
        years = sorted(path.keys())
        prices = [path[y] for y in years]
        color = MPL_SCENARIO_COLORS.get(sid, MPL_BLUE)
        label = SCENARIOS[sid]["name"]
        ax.plot(years, prices, color=color, linewidth=2.5, marker="o", markersize=4, label=label)
    ax.set_xlabel("Year", fontsize=10, color=MPL_DARK)
    ax.set_ylabel("Carbon Price ($/tCO2e)", fontsize=10, color=MPL_DARK)
    ax.set_title("NGFS 시나리오별 탄소가격 경로 (2024-2050)", fontsize=13,
                 fontweight="bold", color=MPL_NAVY, pad=12)
    ax.legend(fontsize=9, framealpha=0.9, loc="upper left")
    ax.grid(True, alpha=0.3)
    fig.tight_layout()
    return fig


def chart_physical_risk_distribution(physical_data):
    """Stacked bar chart of EAL by hazard type for top 8 facilities."""
    fac_results = sorted(physical_data["facilities"],
                         key=lambda x: x["total_expected_annual_loss"], reverse=True)[:8]
    fig, ax = _mpl_setup(figsize=(10, 4.5))

    hazard_types = ["flood", "typhoon", "heatwave", "drought", "sea_level_rise"]
    hazard_colors = [MPL_BLUE, MPL_RED, MPL_ORANGE, MPL_YELLOW, MPL_TEAL]
    fac_names = [f["facility_name"][:8] for f in fac_results]
    x = np.arange(len(fac_names))
    bar_width = 0.6

    bottom = np.zeros(len(fac_results))
    for htype, hcolor in zip(hazard_types, hazard_colors):
        vals = []
        for fac in fac_results:
            hloss = 0
            for h in fac["hazards"]:
                if h["hazard_type"] == htype:
                    hloss = h["potential_loss"] / 1e6
            vals.append(hloss)
        ax.bar(x, vals, bar_width, bottom=bottom, color=hcolor,
               label=HAZARD_NAMES_KR.get(htype, htype), edgecolor="white", linewidth=0.5)
        bottom += np.array(vals)

    ax.set_xticks(x)
    ax.set_xticklabels(fac_names, fontsize=8, rotation=30, ha="right")
    ax.set_ylabel("EAL ($M)", fontsize=10, color=MPL_DARK)
    ax.set_title("시설별 물리적 리스크 분포 (상위 8개)", fontsize=13,
                 fontweight="bold", color=MPL_NAVY, pad=12)
    ax.legend(fontsize=8, ncol=5, loc="upper right", framealpha=0.9)
    fig.tight_layout()
    return fig


def chart_esg_radar(esg_data_dict):
    """Spider/radar chart for 3 ESG frameworks across categories."""
    # Common categories across all frameworks
    common_cats = ["거버넌스", "전략", "리스크 관리", "지표 및 목표"]
    n = len(common_cats)
    angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
    angles += angles[:1]

    if _IS_MAC:
        plt.rcParams["font.family"] = "AppleGothic"
    else:
        plt.rcParams["font.family"] = "Malgun Gothic"
    plt.rcParams["axes.unicode_minus"] = False
    fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
    fig.patch.set_facecolor(MPL_LIGHT_BG)

    fw_colors = {"issb": MPL_BLUE, "tcfd": MPL_ORANGE, "kssb": MPL_GREEN}
    fw_labels = {"issb": "ISSB", "tcfd": "TCFD", "kssb": "KSSB"}

    for fw_id, data in esg_data_dict.items():
        values = []
        for cat_name in common_cats:
            found = False
            for cat in data["categories"]:
                if cat["category"] == cat_name:
                    values.append(cat["score"] / cat["max_score"] * 100)
                    found = True
                    break
            if not found:
                values.append(0)
        values += values[:1]
        color = fw_colors.get(fw_id, MPL_BLUE)
        ax.plot(angles, values, color=color, linewidth=2, label=fw_labels.get(fw_id, fw_id))
        ax.fill(angles, values, color=color, alpha=0.15)

    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(common_cats, fontsize=10)
    ax.set_ylim(0, 100)
    ax.set_yticks([20, 40, 60, 80, 100])
    ax.set_yticklabels(["20", "40", "60", "80", "100"], fontsize=7, color=MPL_GRAY)
    ax.set_title("ESG 프레임워크별 준수 점수", fontsize=13,
                 fontweight="bold", color=MPL_NAVY, pad=20)
    ax.legend(fontsize=9, loc="lower right", bbox_to_anchor=(1.2, 0))
    fig.tight_layout()
    return fig


def chart_kets_allocation_decline():
    """Line chart showing K-ETS free allocation ratio decline by sector."""
    fig, ax = _mpl_setup(figsize=(10, 4.5))
    years = list(range(2024, 2051))
    sectors_to_plot = ["steel", "utilities", "petrochemical", "oil_gas"]
    sector_colors = [MPL_BLUE, MPL_RED, MPL_ORANGE, MPL_GREEN]

    for sec, color in zip(sectors_to_plot, sector_colors):
        params = KETS_ALLOCATION_RATIOS.get(sec)
        if not params:
            continue
        ratios = []
        for y in years:
            r = max(0, params["base_ratio"] - params["annual_tightening"] * (y - params["base_year"]))
            ratios.append(r * 100)
        label = SECTOR_NAMES_KR.get(sec, sec)
        ax.plot(years, ratios, color=color, linewidth=2.5, label=label)

    ax.set_xlabel("Year", fontsize=10, color=MPL_DARK)
    ax.set_ylabel("무상할당 비율 (%)", fontsize=10, color=MPL_DARK)
    ax.set_title("K-ETS 섹터별 무상할당 비율 추이 (2024-2050)", fontsize=13,
                 fontweight="bold", color=MPL_NAVY, pad=12)
    ax.legend(fontsize=9, loc="lower left", framealpha=0.9)
    ax.grid(True, alpha=0.3)
    ax.set_ylim(0, 105)
    fig.tight_layout()
    return fig


def chart_regulatory_timeline():
    """Horizontal timeline with milestone markers."""
    fig, ax = _mpl_setup(figsize=(10, 3.5))

    deadlines = sorted(REGULATORY_DEADLINES.values(), key=lambda d: d["date"])
    colors_cycle = [MPL_RED, MPL_BLUE, MPL_GREEN, MPL_ORANGE, MPL_TEAL]

    for i, dl in enumerate(deadlines):
        dt = datetime.strptime(dl["date"], "%Y-%m-%d")
        color = colors_cycle[i % len(colors_cycle)]
        ax.scatter(dt, 0.5, s=200, color=color, zorder=5, edgecolors="white", linewidths=2)
        ax.annotate(dl["name"],
                    xy=(dt, 0.5),
                    xytext=(0, 20 if i % 2 == 0 else -25),
                    textcoords="offset points",
                    fontsize=8, fontweight="bold", color=MPL_NAVY,
                    ha="center", va="center",
                    arrowprops=dict(arrowstyle="-", color=MPL_GRAY, lw=0.8))

    # Timeline bar
    ax.axhline(y=0.5, color=MPL_GRAY, linewidth=1, alpha=0.5)
    # Today marker
    today = datetime.now()
    ax.axvline(x=today, color=MPL_RED, linewidth=1.5, linestyle="--", alpha=0.6)
    ax.annotate("Today", xy=(today, 0.2), fontsize=8, color=MPL_RED, ha="center")

    ax.set_ylim(0, 1)
    ax.set_yticks([])
    ax.spines["left"].set_visible(False)
    ax.set_title("주요 규제 일정 타임라인", fontsize=13, fontweight="bold",
                 color=MPL_NAVY, pad=12)
    fig.tight_layout()
    return fig


# ══════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS (20 slides)
# ══════════════════════════════════════════════════════════════════════

# ── Part 1: Why Climate Disclosure Matters (Slides 1-4) ──────────

def slide_01_title(prs, facilities):
    """Slide 1: Title slide with headline KPIs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, "0F172A", "1E293B")

    # Decorative accent
    _add_shape(slide, Inches(0), Inches(2.6), Inches(0.1), Inches(2.5), BLUE)

    # Title
    _add_textbox(slide, Inches(1.0), Inches(1.0), Inches(11), Inches(1.0),
                 "Climate Risk Disclosure Advisory",
                 font_size=44, color=WHITE, bold=True)
    # Korean subtitle
    _add_textbox(slide, Inches(1.0), Inches(2.0), Inches(11), Inches(0.6),
                 "기후 리스크 공시 자문 제안서", font_size=24, color=BLUE)
    # Date
    _add_textbox(slide, Inches(1.0), Inches(2.8), Inches(11), Inches(0.4),
                 f"{date.today().isoformat()}  |  Confidential",
                 font_size=14, color=MID_GRAY)

    # 4 headline KPIs
    total_s1 = sum(f["current_emissions_scope1"] for f in facilities)
    total_s2 = sum(f["current_emissions_scope2"] for f in facilities)
    total_assets = sum(f["assets_value"] for f in facilities)
    sectors = set(f["sector"] for f in facilities)

    kpis = [
        ("분석 대상 시설", f"{len(facilities)}개"),
        ("분석 대상 섹터", f"{len(sectors)}개"),
        ("총 배출량 (S1+S2)", f"{(total_s1+total_s2)/1e6:.0f}M tCO2e"),
        ("총 자산 규모", f"${total_assets/1e9:.0f}B"),
    ]
    for i, (label, value) in enumerate(kpis):
        x = Inches(1.0) + i * Inches(3.1)
        _add_kpi_box(slide, x, Inches(4.0), Inches(2.8), Inches(1.1),
                     label, value, bg_color=NAVY, value_color=WHITE)

    # Bottom line
    _add_shape(slide, Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.02), MID_GRAY)
    _add_textbox(slide, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.6),
                 "NGFS 시나리오 기반 전환·물리적 리스크 통합 분석  |  ESG 공시 준비도 평가",
                 font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def slide_02_executive_summary(prs, summary_data, esg_kssb, physical_data):
    """Slide 2: Executive Summary - 4 KPIs + key findings."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Executive Summary", "경영진 요약  |  주요 분석 결과")

    # 4 KPI boxes
    total_npv = summary_data["total_npv"]
    high_risk = summary_data["high_risk_count"]
    esg_score = esg_kssb["overall_score"]
    total_eal = sum(f["total_expected_annual_loss"] for f in physical_data["facilities"])

    kpis = [
        ("전환 리스크 NPV", f"-{_fmt_b(total_npv)}", LIGHT_RED, RED),
        ("고위험 시설 수", f"{high_risk}개", LIGHT_ORANGE, ORANGE),
        ("KSSB 준수 점수", f"{esg_score:.0f}%", LIGHT_BLUE, BLUE),
        ("물리적 리스크 EAL", f"${total_eal/1e6:.1f}M/yr", LIGHT_GREEN, TEAL),
    ]
    for i, (label, value, bg, vc) in enumerate(kpis):
        x = LEFT_MARGIN + i * Inches(3.0)
        _add_kpi_box(slide, x, Inches(1.6), Inches(2.7), Inches(1.1), label, value,
                     bg_color=bg, value_color=vc)

    # Key findings
    _add_textbox(slide, LEFT_MARGIN, Inches(3.1), CONTENT_WIDTH, Inches(0.5),
                 "Key Findings  |  주요 발견사항", font_size=18, color=NAVY, bold=True)

    findings = [
        f"Net Zero 2050 시나리오 기준, 포트폴리오 전체 NPV 영향 -{_fmt_b(total_npv)} (자산 대비 {abs(total_npv)/sum(f['assets_value'] for f in get_all_facilities())*100:.1f}%)",
        f"전체 {summary_data['total_facilities']}개 시설 중 {high_risk}개 시설이 고위험(High Risk) 등급",
        f"물리적 리스크 연간 예상 손실(EAL) ${total_eal/1e6:.1f}M — 홍수·태풍이 주요 위험 요인",
        f"KSSB 공시 준비도 {esg_score:.0f}% — 2025년 의무 공시 기한 대비 추가 준비 필요",
    ]
    _add_bullet_list(slide, LEFT_MARGIN, Inches(3.6), CONTENT_WIDTH, Inches(3.0),
                     findings, font_size=13)

    _add_footer(slide, 2)


def slide_03_regulatory_landscape(prs):
    """Slide 3: Regulatory Landscape - 3 framework cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Regulatory Landscape", "규제 환경  |  주요 기후 공시 프레임워크")

    # 3 cards: TCFD/ISSB, KSSB, EU CBAM
    cards = [
        ("TCFD / ISSB", RED, [
            "글로벌 기후 공시 표준",
            "IFRS S1/S2: 2024.01 발효",
            "거버넌스·전략·리스크관리·지표",
            "시나리오 분석 필수 요구",
            "Scope 1/2/3 배출량 공시",
        ]),
        ("KSSB 한국 공시기준", BLUE, [
            "한국 지속가능성기준위원회",
            "2025: 자산 2조원+ 의무 공시",
            "2027: 전 상장사 확대 적용",
            "ISSB 기반 + 한국 특화 요건",
            "K-ETS 연계 배출권 정보 공시",
        ]),
        ("EU CBAM", GREEN, [
            "EU 탄소국경조정제도",
            "2026.01 본격 시행",
            "철강·시멘트·알루미늄 대상",
            "내재 탄소 비용 과세",
            "한국 수출기업 직접 영향",
        ]),
    ]

    for i, (title, accent, lines) in enumerate(cards):
        x = LEFT_MARGIN + i * Inches(4.0)
        _add_card(slide, x, Inches(1.6), Inches(3.7), Inches(3.8),
                  title, lines, accent_color=accent)

    # Bottom note
    _add_rounded_rect(slide, LEFT_MARGIN, Inches(5.8), CONTENT_WIDTH, Inches(0.8), LIGHT_BLUE, shadow=False)
    _add_textbox(slide, LEFT_MARGIN + Inches(0.3), Inches(5.9), CONTENT_WIDTH - Inches(0.6), Inches(0.6),
                 "규제 미준수 시: 과징금, 투자자 신뢰 하락, ESG 평가 등급 하락, 수출 경쟁력 저하",
                 font_size=13, color=NAVY, bold=True)

    _add_footer(slide, 3)


def slide_04_why_now(prs):
    """Slide 4: Why Now? - Timeline chart + urgency points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Why Now?", "왜 지금인가  |  규제 타임라인과 시급성")

    # Regulatory timeline chart
    fig = chart_regulatory_timeline()
    _embed_chart(slide, fig, LEFT_MARGIN, Inches(1.5), Inches(11.5), Inches(2.8))

    # 3 urgency cards
    urgency = [
        ("규제 압력 가속", RED,
         "KSSB 의무 공시 2025년 시행. 준비 기간 1년 미만."),
        ("투자자 요구 증대", ORANGE,
         "글로벌 기관투자자 73%가 기후 리스크 공시를 투자 결정에 반영 (PwC 2024)."),
        ("선제적 대응 = 경쟁 우위", GREEN,
         "조기 공시 기업은 자본비용 30-50bp 절감 효과 (Bolton & Kacperczyk 2021)."),
    ]
    for i, (title, accent, desc) in enumerate(urgency):
        x = LEFT_MARGIN + i * Inches(4.0)
        _add_card(slide, x, Inches(4.6), Inches(3.7), Inches(2.0),
                  title, [desc], accent_color=accent)

    _add_footer(slide, 4)


# ── Part 2: Your Risk Profile (Slides 5-12) ──────────────────────

def slide_05_section_divider(prs):
    """Slide 5: Section divider - Your Climate Risk Profile."""
    slide = _add_section_divider(prs, 2, "Your Climate Risk Profile",
                                  "귀사의 기후 리스크 프로파일")
    return slide


def slide_06_portfolio_overview(prs, facilities):
    """Slide 6: Full facility table + sector distribution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Portfolio Overview", "포트폴리오 개요  |  분석 대상 시설 현황")

    # Sector distribution KPIs
    sector_counts = Counter(f["sector"] for f in facilities)
    top_sectors = sector_counts.most_common(4)
    for i, (sec, cnt) in enumerate(top_sectors):
        x = LEFT_MARGIN + i * Inches(3.0)
        _add_kpi_box(slide, x, Inches(1.5), Inches(2.7), Inches(0.9),
                     SECTOR_NAMES_KR.get(sec, sec), f"{cnt}개 시설",
                     bg_color=LIGHT_BLUE, value_color=BLUE)

    # Facility table (compact)
    headers = ["시설명", "섹터", "배출량 (S1+S2)", "자산 규모", "매출액"]
    rows = []
    for f in facilities:
        emissions = (f["current_emissions_scope1"] + f["current_emissions_scope2"]) / 1e6
        rows.append([
            f["name"],
            SECTOR_NAMES_KR.get(f["sector"], f["sector"]),
            f"{emissions:.1f}M tCO2e",
            f"${f['assets_value']/1e9:.1f}B",
            f"${f['annual_revenue']/1e9:.1f}B",
        ])

    _add_table(slide, LEFT_MARGIN, Inches(2.7), CONTENT_WIDTH, Inches(4.2),
               headers, rows,
               col_widths=[Inches(2.8), Inches(2.0), Inches(2.4), Inches(2.2), Inches(2.3)])

    _add_footer(slide, 6)


def slide_07_transition_risk_summary(prs, comparison_data):
    """Slide 7: 4-scenario NPV comparison bar chart + risk distribution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Transition Risk Summary",
                   "전환 리스크 요약  |  NGFS 4대 시나리오 비교")

    # NPV bar chart
    fig = chart_scenario_npv_comparison(comparison_data)
    _embed_chart(slide, fig, LEFT_MARGIN, Inches(1.5), Inches(7.5), Inches(3.8))

    # Risk distribution boxes
    _add_textbox(slide, Inches(9.0), Inches(1.6), Inches(3.5), Inches(0.4),
                 "리스크 분포", font_size=16, color=NAVY, bold=True)

    y_off = Inches(2.2)
    risk_dist = comparison_data.get("risk_distribution", {})
    for sid in ["net_zero_2050", "below_2c", "delayed_transition", "current_policies"]:
        dist = risk_dist.get(sid, {})
        high = dist.get("high", 0)
        med = dist.get("medium", 0)
        low = dist.get("low", 0)
        color = SCENARIO_COLORS.get(sid, BLUE)
        name = SCENARIO_NAMES_KR.get(sid, sid)

        _add_rounded_rect(slide, Inches(9.0), y_off, Inches(3.5), Inches(0.8),
                          LIGHT_GRAY, shadow=False)
        _add_shape(slide, Inches(9.0), y_off, Inches(0.08), Inches(0.8), color)
        _add_textbox(slide, Inches(9.2), y_off + Inches(0.05), Inches(3.0), Inches(0.3),
                     name, font_size=11, color=NAVY, bold=True)
        _add_rich_text(slide, Inches(9.2), y_off + Inches(0.35), Inches(3.0), Inches(0.35),
                       [
                           ("H:", 10, RED, True), (f" {high}", 10, DARK_GRAY, False),
                           ("  M:", 10, ORANGE, True), (f" {med}", 10, DARK_GRAY, False),
                           ("  L:", 10, GREEN, True), (f" {low}", 10, DARK_GRAY, False),
                       ])
        y_off += Inches(0.9)

    _add_footer(slide, 7)


def slide_08_cost_breakdown(prs, summary_data, analysis_data):
    """Slide 8: Cost waterfall + Top 5 facilities table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Transition Cost Breakdown",
                   "전환 비용 분석  |  비용 구성 및 고위험 시설")

    # Cost waterfall chart
    fig = chart_cost_waterfall(summary_data)
    _embed_chart(slide, fig, LEFT_MARGIN, Inches(1.5), Inches(6.5), Inches(3.5))

    # Top 5 facilities
    _add_textbox(slide, Inches(8.0), Inches(1.5), Inches(4.5), Inches(0.4),
                 "Top 5 고위험 시설", font_size=16, color=NAVY, bold=True)

    risk_colors = {"High": RED, "Medium": ORANGE, "Low": GREEN}
    facs = sorted(analysis_data["facilities"], key=lambda x: x["delta_npv"])
    top5_headers = ["시설명", "Delta NPV", "리스크"]
    top5_rows = []
    for f in facs[:5]:
        level = f["risk_level"]
        top5_rows.append([
            f["facility_name"],
            f"-{_fmt_b(f['delta_npv'])}",
            (level, risk_colors.get(level, BLACK), True),
        ])

    _add_table(slide, Inches(8.0), Inches(2.0), Inches(4.5), Inches(2.2),
               top5_headers, top5_rows,
               col_widths=[Inches(1.8), Inches(1.4), Inches(1.3)])

    # Cost summary bar
    breakdown = summary_data["cost_breakdown"]
    _add_rounded_rect(slide, LEFT_MARGIN, Inches(5.3), CONTENT_WIDTH, Inches(0.9), NAVY, shadow=True)
    _add_rich_text(slide, Inches(1.0), Inches(5.4), Inches(11.0), Inches(0.5),
                   [
                       ("총 NPV: ", 14, WHITE, True),
                       (f"-{_fmt_b(summary_data['total_npv'])}", 14, LIGHT_RED, True),
                       ("    탄소: ", 12, MID_GRAY, False),
                       (f"{_fmt_b(breakdown['carbon_cost'])}", 12, LIGHT_BLUE, True),
                       ("    에너지: ", 12, MID_GRAY, False),
                       (f"{_fmt_b(breakdown['energy_cost_increase'])}", 12, LIGHT_BLUE, True),
                       ("    매출: ", 12, MID_GRAY, False),
                       (f"{_fmt_b(breakdown['revenue_impact'])}", 12, LIGHT_BLUE, True),
                   ])

    _add_footer(slide, 8)


def slide_09_carbon_price(prs):
    """Slide 9: Carbon price pathway chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Carbon Price Pathways",
                   "탄소가격 경로  |  NGFS 시나리오별 전망 (2024-2050)")

    fig = chart_carbon_price_paths()
    _embed_chart(slide, fig, LEFT_MARGIN, Inches(1.5), Inches(11.5), Inches(4.5))

    # Key takeaway
    _add_rounded_rect(slide, LEFT_MARGIN, Inches(6.2), CONTENT_WIDTH, Inches(0.7),
                      LIGHT_BLUE, shadow=False)
    nz_2030 = NGFS_CARBON_PRICE_PATHS["net_zero_2050"][2030]
    nz_2050 = NGFS_CARBON_PRICE_PATHS["net_zero_2050"][2050]
    _add_textbox(slide, LEFT_MARGIN + Inches(0.3), Inches(6.3), CONTENT_WIDTH - Inches(0.6), Inches(0.4),
                 f"Net Zero 2050: 2030년 ${nz_2030}/tCO2e → 2050년 ${nz_2050}/tCO2e  |  "
                 "탄소 집약적 산업에 대한 비용 부담 급격히 증가",
                 font_size=12, color=NAVY, bold=True)

    _add_footer(slide, 9)


def slide_10_physical_risk_summary(prs, physical_data):
    """Slide 10: Physical risk EAL distribution chart + summary KPIs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Physical Risk Summary",
                   "물리적 리스크 요약  |  기후재해 유형별 시설 취약도")

    # EAL stacked bar chart
    fig = chart_physical_risk_distribution(physical_data)
    _embed_chart(slide, fig, LEFT_MARGIN, Inches(1.5), Inches(7.5), Inches(3.8))

    # Summary KPIs
    total_eal = sum(f["total_expected_annual_loss"] for f in physical_data["facilities"])
    summary = physical_data["overall_risk_summary"]
    high_count = summary.get("High", 0)
    med_count = summary.get("Medium", 0)

    _add_textbox(slide, Inches(9.0), Inches(1.5), Inches(3.5), Inches(0.4),
                 "리스크 요약", font_size=16, color=NAVY, bold=True)

    kpi_data = [
        ("총 EAL (연간)", f"${total_eal/1e6:.1f}M", LIGHT_BLUE, BLUE),
        ("고위험 시설", f"{high_count}개", LIGHT_RED, RED),
        ("중위험 시설", f"{med_count}개", LIGHT_ORANGE, ORANGE),
        ("분석 시나리오", physical_data.get("scenario", "N/A"), LIGHT_GRAY, DARK_GRAY),
    ]
    for i, (label, value, bg, vc) in enumerate(kpi_data):
        y = Inches(2.1) + i * Inches(1.0)
        _add_kpi_box(slide, Inches(9.0), y, Inches(3.5), Inches(0.85), label, value,
                     bg_color=bg, value_color=vc)

    _add_footer(slide, 10)


def slide_11_physical_risk_detail(prs, physical_data):
    """Slide 11: Top 5 EAL table + 5 hazard mini-cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Physical Risk Detail",
                   "물리적 리스크 상세  |  시설별 EAL 및 재해 유형")

    # 5 hazard mini-cards
    hazards_info = [
        ("홍수", "집중호우/하천범람", BLUE),
        ("태풍", "강풍/폭풍해일", RED),
        ("폭염", "냉각효율/생산성", ORANGE),
        ("가뭄", "용수부족/냉각수", YELLOW),
        ("해수면", "연안시설 침수", TEAL),
    ]
    for i, (title, desc, color) in enumerate(hazards_info):
        x = Inches(0.5) + i * Inches(2.5)
        _add_rounded_rect(slide, x, Inches(1.5), Inches(2.2), Inches(1.0), LIGHT_BLUE, shadow=True)
        _add_shape(slide, x, Inches(1.5), Inches(2.2), Inches(0.05), color)
        _add_textbox(slide, x + Inches(0.1), Inches(1.6), Inches(2.0), Inches(0.35),
                     title, font_size=14, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.1), Inches(1.95), Inches(2.0), Inches(0.35),
                     desc, font_size=10, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # EAL Top 5 table
    _add_textbox(slide, LEFT_MARGIN, Inches(2.8), Inches(11), Inches(0.5),
                 "시설별 연간예상손실(EAL) Top 5", font_size=18, color=NAVY, bold=True)

    fac_results = sorted(physical_data["facilities"],
                         key=lambda x: x["total_expected_annual_loss"], reverse=True)

    headers = ["순위", "시설명", "위치", "위험등급", "EAL ($M)", "주요 재해"]
    rows = []
    for i, fac in enumerate(fac_results[:5]):
        top_hazard = max(fac["hazards"], key=lambda h: h["potential_loss"])
        level = fac["overall_risk_level"]
        risk_colors = {"High": RED, "Medium": ORANGE, "Low": GREEN}
        rows.append([
            str(i + 1),
            fac["facility_name"],
            fac["location"],
            (level, risk_colors.get(level, BLACK), True),
            f"${fac['total_expected_annual_loss']/1e6:.1f}",
            HAZARD_NAMES_KR.get(top_hazard["hazard_type"], top_hazard["hazard_type"]),
        ])

    _add_table(slide, LEFT_MARGIN, Inches(3.3), CONTENT_WIDTH, Inches(2.2),
               headers, rows,
               col_widths=[Inches(0.8), Inches(2.5), Inches(2.0), Inches(1.8), Inches(2.2), Inches(2.4)])

    # Summary bar
    total_eal = sum(f["total_expected_annual_loss"] for f in physical_data["facilities"])
    summary = physical_data["overall_risk_summary"]
    _add_rounded_rect(slide, LEFT_MARGIN, Inches(5.8), CONTENT_WIDTH, Inches(0.8), NAVY, shadow=True)
    _add_rich_text(slide, Inches(1.0), Inches(5.9), Inches(11.0), Inches(0.5),
                   [
                       ("총 EAL: ", 14, WHITE, True),
                       (f"${total_eal/1e6:.1f}M", 14, LIGHT_BLUE, True),
                       ("    ", 14, WHITE, False),
                       (f"High {summary.get('High', 0)}개", 14, LIGHT_RED, True),
                       ("  |  ", 14, MID_GRAY, False),
                       (f"Medium {summary.get('Medium', 0)}개", 14, LIGHT_ORANGE, True),
                       ("  |  ", 14, MID_GRAY, False),
                       (f"Low {summary.get('Low', 0)}개", 14, LIGHT_GREEN, True),
                   ])

    _add_footer(slide, 11)


def slide_12_kets_impact(prs, comparison_data):
    """Slide 12: K-ETS free allocation decline + Global vs K-ETS comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "K-ETS Impact Analysis",
                   "K-ETS 영향 분석  |  무상할당 축소 전망 및 글로벌 비교")

    # K-ETS allocation decline chart
    fig = chart_kets_allocation_decline()
    _embed_chart(slide, fig, LEFT_MARGIN, Inches(1.5), Inches(7.0), Inches(3.8))

    # K-ETS info cards
    _add_textbox(slide, Inches(8.5), Inches(1.5), Inches(4.0), Inches(0.4),
                 "K-ETS 주요 변동", font_size=16, color=NAVY, bold=True)

    kets_items = [
        ("Phase 4 (2026~)", "강화된 할당 기준\n연간 1-2%p 무상할당 축소", RED),
        ("EITE 업종", "철강·시멘트 97% → 점진 감소\n수출 경쟁력 보호 한시적", BLUE),
        ("발전 부문", "90% 무상할당 시작\n연간 1.5%p 가속 축소", ORANGE),
        ("비용 영향", "초과 배출분만 과금\n글로벌 대비 비용 절감 효과", GREEN),
    ]
    for i, (title, desc, accent) in enumerate(kets_items):
        y = Inches(2.0) + i * Inches(1.2)
        _add_rounded_rect(slide, Inches(8.5), y, Inches(4.0), Inches(1.05),
                          LIGHT_GRAY, shadow=False)
        _add_shape(slide, Inches(8.5), y, Inches(0.06), Inches(1.05), accent)
        _add_textbox(slide, Inches(8.7), y + Inches(0.05), Inches(3.6), Inches(0.3),
                     title, font_size=12, color=NAVY, bold=True)
        _add_textbox(slide, Inches(8.7), y + Inches(0.35), Inches(3.6), Inches(0.6),
                     desc, font_size=10, color=DARK_GRAY)

    _add_footer(slide, 12)


# ── Part 3: ESG Readiness (Slides 13-16) ─────────────────────────

def slide_13_section_divider(prs):
    """Slide 13: Section divider - ESG Disclosure Readiness."""
    _add_section_divider(prs, 3, "ESG Disclosure Readiness",
                          "ESG 공시 준비도 평가")


def slide_14_esg_compliance(prs, esg_data_dict):
    """Slide 14: ESG radar chart + overall score + maturity."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "ESG Compliance Score",
                   "ESG 준수 점수  |  3대 프레임워크 평가 결과")

    # Radar chart
    fig = chart_esg_radar(esg_data_dict)
    _embed_chart(slide, fig, LEFT_MARGIN, Inches(1.5), Inches(5.5), Inches(5.0))

    # Framework scores
    _add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4),
                 "프레임워크별 점수", font_size=16, color=NAVY, bold=True)

    fw_colors = {"issb": BLUE, "tcfd": ORANGE, "kssb": GREEN}
    fw_names = {"issb": "ISSB (IFRS S2)", "tcfd": "TCFD", "kssb": "KSSB"}

    y_off = Inches(2.1)
    for fw_id, data in esg_data_dict.items():
        color = fw_colors.get(fw_id, BLUE)
        name = fw_names.get(fw_id, fw_id)
        score = data["overall_score"]
        compliance = data["compliance_level"]

        _add_rounded_rect(slide, Inches(7.0), y_off, Inches(5.5), Inches(1.2),
                          LIGHT_GRAY, shadow=False)
        _add_shape(slide, Inches(7.0), y_off, Inches(0.08), Inches(1.2), color)
        _add_textbox(slide, Inches(7.3), y_off + Inches(0.1), Inches(5.0), Inches(0.3),
                     name, font_size=14, color=NAVY, bold=True)
        _add_rich_text(slide, Inches(7.3), y_off + Inches(0.45), Inches(5.0), Inches(0.3),
                       [
                           ("점수: ", 12, DARK_GRAY, False),
                           (f"{score:.0f}%", 16, color, True),
                           ("    준수 수준: ", 12, DARK_GRAY, False),
                           (compliance, 12, NAVY, True),
                       ])

        # Maturity level
        maturity = data.get("maturity_level")
        if maturity:
            _add_textbox(slide, Inches(7.3), y_off + Inches(0.8), Inches(5.0), Inches(0.3),
                         f"성숙도: Lv.{maturity['level']} {maturity['name']}",
                         font_size=11, color=MID_GRAY)
        y_off += Inches(1.4)

    _add_footer(slide, 14)


def slide_15_gap_analysis(prs, esg_kssb):
    """Slide 15: Gap analysis priority table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Gap Analysis", "갭 분석  |  KSSB 준수를 위한 우선순위 과제")

    gap_analysis = esg_kssb.get("gap_analysis", [])

    if gap_analysis:
        headers = ["카테고리", "현재 점수", "목표 점수", "갭", "우선순위", "난이도"]
        rows = []
        effort_colors = {"High": RED, "Medium": ORANGE, "Low": GREEN}
        for gap in gap_analysis:
            effort = gap.get("effort", "Medium")
            rows.append([
                gap["category"],
                f"{gap['current_score']:.0f}",
                f"{gap['target_score']:.0f}",
                f"{gap['gap']:.0f}",
                f"{gap['priority_score']:.1f}",
                (effort, effort_colors.get(effort, BLACK), True),
            ])

        _add_table(slide, LEFT_MARGIN, Inches(1.6), CONTENT_WIDTH, Inches(2.5),
                   headers, rows,
                   col_widths=[Inches(2.5), Inches(1.5), Inches(1.5), Inches(1.5),
                               Inches(2.0), Inches(2.7)])

    # Recommended actions
    _add_textbox(slide, LEFT_MARGIN, Inches(4.4), CONTENT_WIDTH, Inches(0.5),
                 "권장 조치 사항", font_size=18, color=NAVY, bold=True)

    recommendations = esg_kssb.get("recommendations", [])
    display_recs = recommendations[:5] if recommendations else [
        "기후 리스크 거버넌스 체계 수립 (이사회 ESG 위원회 설치)",
        "시나리오 분석 기반 전략적 기후 대응 계획 수립",
        "Scope 1/2/3 배출량 산정 및 검증 체계 구축",
        "KSSB 공시 양식에 맞춘 데이터 수집 프로세스 정비",
        "내부 감사 및 외부 인증 체계 마련",
    ]
    _add_bullet_list(slide, LEFT_MARGIN, Inches(4.9), CONTENT_WIDTH, Inches(2.0),
                     display_recs, font_size=12)

    _add_footer(slide, 15)


def slide_16_regulatory_timeline(prs):
    """Slide 16: Full regulatory timeline + deadline table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Regulatory Timeline",
                   "규제 일정  |  주요 기후 공시 마감일 및 D-Day")

    # Deadline table with D-day calculation
    today = date.today()
    headers = ["규제", "시행일", "설명", "D-Day", "상태"]
    rows = []
    for key, dl in sorted(REGULATORY_DEADLINES.items(), key=lambda x: x[1]["date"]):
        dl_date = datetime.strptime(dl["date"], "%Y-%m-%d").date()
        d_day = (dl_date - today).days
        if d_day < 0:
            status_text = "시행 중"
            status_color = GREEN
        elif d_day < 365:
            status_text = "긴급"
            status_color = RED
        else:
            status_text = "준비 중"
            status_color = ORANGE

        d_day_str = f"D{d_day:+d}" if d_day != 0 else "D-Day"
        rows.append([
            dl["name"],
            dl["date"],
            dl["description"][:30] + "..." if len(dl["description"]) > 30 else dl["description"],
            d_day_str,
            (status_text, status_color, True),
        ])

    _add_table(slide, LEFT_MARGIN, Inches(1.6), CONTENT_WIDTH, Inches(2.8),
               headers, rows,
               col_widths=[Inches(2.2), Inches(1.8), Inches(3.5), Inches(1.5), Inches(2.7)])

    # Action items
    _add_textbox(slide, LEFT_MARGIN, Inches(4.8), CONTENT_WIDTH, Inches(0.5),
                 "대응 전략", font_size=18, color=NAVY, bold=True)

    actions = [
        "즉시: KSSB 공시 준비 태스크포스 구성 및 갭 분석 착수",
        "6개월 내: Scope 1/2 배출량 산정 체계 구축 및 내부 검증",
        "12개월 내: 시나리오 분석 완료 및 KSSB 양식 초안 작성",
        "18개월 내: 외부 인증 및 공시 보고서 최종 발행",
    ]
    _add_bullet_list(slide, LEFT_MARGIN, Inches(5.3), CONTENT_WIDTH, Inches(1.5),
                     actions, font_size=13)

    _add_footer(slide, 16)


# ── Part 4: Our Solution (Slides 17-20) ──────────────────────────

def slide_17_section_divider(prs):
    """Slide 17: Section divider - Our Solution."""
    _add_section_divider(prs, 4, "Our Solution",
                          "Climate Risk Platform을 활용한 통합 솔루션")


def slide_18_platform_capabilities(prs):
    """Slide 18: Platform capabilities - 3 columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Platform Capabilities",
                   "플랫폼 역량  |  통합 기후 리스크 분석 솔루션")

    # 3 capability columns
    capabilities = [
        ("전환 리스크 분석", BLUE, [
            "NGFS 4대 시나리오 기반",
            "S-curve 배출 경로 모델링",
            "섹터별 기술 감축 스택",
            "탄소가격 경로 시뮬레이션",
            "K-ETS 무상할당 반영",
            "좌초자산 분석 (Carbon Tracker)",
            "시나리오 NPV 비교 분석",
        ]),
        ("물리적 리스크 분석", RED, [
            "5대 재해 유형 분석",
            "Gumbel 극치통계 모델",
            "KMA 30년 관측 데이터",
            "Open-Meteo API 연동",
            "USACE 피해함수 적용",
            "EAL(연간예상손실) 산정",
            "복합 리스크 분석",
        ]),
        ("ESG 공시 지원", GREEN, [
            "KSSB/ISSB/TCFD 프레임워크",
            "데이터 기반 준수도 평가",
            "갭 분석 및 우선순위 도출",
            "성숙도 모델 (5단계)",
            "규제 일정 자동 추적",
            "공시 데이터 자동 생성",
            "체크리스트 기반 점검",
        ]),
    ]

    for i, (title, accent, items) in enumerate(capabilities):
        x = LEFT_MARGIN + i * Inches(4.0)
        _add_card(slide, x, Inches(1.5), Inches(3.7), Inches(4.8),
                  title, items, accent_color=accent)

    # Bottom tech stack
    _add_rounded_rect(slide, LEFT_MARGIN, Inches(6.5), CONTENT_WIDTH, Inches(0.5),
                      NAVY, shadow=False)
    _add_textbox(slide, LEFT_MARGIN + Inches(0.2), Inches(6.55), CONTENT_WIDTH - Inches(0.4), Inches(0.4),
                 "Tech Stack:  FastAPI  |  Next.js  |  Streamlit  |  Python  |  NGFS  |  Open-Meteo  |  python-pptx",
                 font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)

    _add_footer(slide, 18)


def slide_19_engagement_roadmap(prs):
    """Slide 19: 3-phase engagement roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Engagement Roadmap",
                   "프로젝트 로드맵  |  3단계 실행 계획")

    phases = [
        ("Phase 1", "Assessment", "현황 진단 (4-6주)", BLUE, [
            "포트폴리오 데이터 수집 및 정리",
            "전환·물리적 리스크 초기 분석",
            "ESG 공시 갭 분석 수행",
            "현행 거버넌스 체계 점검",
            "경영진 보고 (중간 결과)",
        ]),
        ("Phase 2", "Strategy", "전략 수립 (6-8주)", ORANGE, [
            "시나리오별 심층 리스크 분석",
            "탈탄소 전환 전략 개발",
            "재무적 영향 정량화",
            "K-ETS 대응 전략 수립",
            "이사회 보고 및 승인",
        ]),
        ("Phase 3", "Reporting", "공시 작성 (4-6주)", GREEN, [
            "KSSB 공시 양식 작성",
            "데이터 검증 및 인증 준비",
            "경영진 리뷰 및 최종 승인",
            "규제 기관 제출",
            "지속적 모니터링 체계 구축",
        ]),
    ]

    for i, (phase_num, phase_name, timeline, accent, items) in enumerate(phases):
        x = LEFT_MARGIN + i * Inches(4.0)
        # Phase card
        _add_rounded_rect(slide, x, Inches(1.5), Inches(3.7), Inches(5.0), WHITE, shadow=True)
        _add_shape(slide, x, Inches(1.5), Inches(3.7), Inches(0.06), accent)

        # Phase header
        _add_rounded_rect(slide, x + Inches(0.15), Inches(1.7), Inches(3.4), Inches(0.9),
                          accent, shadow=False)
        _add_textbox(slide, x + Inches(0.3), Inches(1.75), Inches(3.1), Inches(0.35),
                     f"{phase_num}: {phase_name}", font_size=18, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.3), Inches(2.1), Inches(3.1), Inches(0.35),
                     timeline, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

        # Items
        _add_bullet_list(slide, x + Inches(0.2), Inches(2.8), Inches(3.3), Inches(3.0),
                         items, font_size=11, bullet_color=accent)

    # Connector arrows
    for i in range(2):
        x = LEFT_MARGIN + (i + 1) * Inches(4.0) - Inches(0.15)
        _add_textbox(slide, x, Inches(3.8), Inches(0.3), Inches(0.4),
                     "\u25B6", font_size=20, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    _add_footer(slide, 19)


def slide_20_contact(prs):
    """Slide 20: Contact / CTA."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, "0F172A", "1E293B")

    # Accent
    _add_shape(slide, Inches(0), Inches(2.5), Inches(0.1), Inches(2.5), BLUE)

    # Main CTA
    _add_textbox(slide, Inches(1.2), Inches(1.5), Inches(11), Inches(1.0),
                 "Ready to Start?", font_size=48, color=WHITE, bold=True)
    _add_textbox(slide, Inches(1.2), Inches(2.6), Inches(11), Inches(0.7),
                 "기후 리스크 공시, 지금 시작하세요",
                 font_size=24, color=BLUE)

    # Contact cards
    _add_shape(slide, Inches(1.2), Inches(3.8), Inches(10), Inches(0.02), MID_GRAY)

    contact_items = [
        ("다음 단계", "무료 포트폴리오 진단 (2주) → Phase 1 Assessment 시작"),
        ("문의", "[담당자 이름]  |  [email@company.com]  |  [02-XXXX-XXXX]"),
        ("웹사이트", "[Climate Risk Platform URL]"),
    ]

    y_off = Inches(4.2)
    for title, desc in contact_items:
        _add_textbox(slide, Inches(1.2), y_off, Inches(3.0), Inches(0.4),
                     title, font_size=16, color=BLUE, bold=True)
        _add_textbox(slide, Inches(4.5), y_off, Inches(7.0), Inches(0.4),
                     desc, font_size=14, color=WHITE)
        y_off += Inches(0.6)

    # Bottom
    _add_shape(slide, Inches(1.2), Inches(6.2), Inches(10), Inches(0.02), MID_GRAY)
    _add_textbox(slide, Inches(1.2), Inches(6.4), Inches(10), Inches(0.4),
                 "Climate Risk Advisory  |  Confidential  |  " + date.today().isoformat(),
                 font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# MAIN ORCHESTRATOR
# ══════════════════════════════════════════════════════════════════════

def generate_proposal(
    scenario_id: str = "net_zero_2050",
    pricing_regime: str = "kets",
    year: int = 2030,
    output_path: str | None = None,
) -> str:
    """Generate the full 20-slide client proposal PPT."""
    print("  [1/5] Loading data...")
    facilities = get_all_facilities()
    analysis = analyse_scenario(scenario_id, pricing_regime=pricing_regime)
    comparison = compare_scenarios(pricing_regime=pricing_regime)
    summary = get_summary(scenario_id, pricing_regime=pricing_regime)
    physical = assess_physical_risk(scenario_id=scenario_id, year=year)

    print("  [2/5] Running ESG assessments...")
    esg_issb = assess_framework("issb")
    esg_tcfd = assess_framework("tcfd")
    esg_kssb = assess_framework("kssb")
    esg_data_dict = {"issb": esg_issb, "tcfd": esg_tcfd, "kssb": esg_kssb}

    print("  [3/5] Creating presentation...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Part 1: Why Climate Disclosure Matters (Slides 1-4)
    slide_01_title(prs, facilities)
    slide_02_executive_summary(prs, summary, esg_kssb, physical)
    slide_03_regulatory_landscape(prs)
    slide_04_why_now(prs)

    print("  [4/5] Building risk profile slides...")
    # Part 2: Your Risk Profile (Slides 5-12)
    slide_05_section_divider(prs)
    slide_06_portfolio_overview(prs, facilities)
    slide_07_transition_risk_summary(prs, comparison)
    slide_08_cost_breakdown(prs, summary, analysis)
    slide_09_carbon_price(prs)
    slide_10_physical_risk_summary(prs, physical)
    slide_11_physical_risk_detail(prs, physical)
    slide_12_kets_impact(prs, comparison)

    # Part 3: ESG Readiness (Slides 13-16)
    slide_13_section_divider(prs)
    slide_14_esg_compliance(prs, esg_data_dict)
    slide_15_gap_analysis(prs, esg_kssb)
    slide_16_regulatory_timeline(prs)

    # Part 4: Our Solution (Slides 17-20)
    slide_17_section_divider(prs)
    slide_18_platform_capabilities(prs)
    slide_19_engagement_roadmap(prs)
    slide_20_contact(prs)

    print("  [5/5] Saving...")
    if output_path is None:
        output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs")
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, "climate_risk_proposal.pptx")

    prs.save(output_path)
    return os.path.abspath(output_path)


def main():
    parser = argparse.ArgumentParser(
        description="Generate a 20-slide climate risk disclosure advisory proposal PPT."
    )
    parser.add_argument(
        "--scenario", default="net_zero_2050",
        choices=list(SCENARIOS.keys()),
        help="Primary NGFS scenario (default: net_zero_2050)",
    )
    parser.add_argument(
        "--pricing", default="kets",
        choices=["global", "kets"],
        help="Carbon pricing regime (default: kets)",
    )
    parser.add_argument(
        "--year", type=int, default=2030,
        help="Physical risk assessment year (default: 2030)",
    )
    parser.add_argument(
        "--output", default=None,
        help="Output file path (default: outputs/climate_risk_proposal.pptx)",
    )
    args = parser.parse_args()

    print("=" * 60)
    print("  Climate Risk Advisory - Client Proposal (20 slides)")
    print("=" * 60)
    print(f"  Scenario: {args.scenario}")
    print(f"  Pricing:  {args.pricing}")
    print(f"  Year:     {args.year}")
    print(f"  Font:     {LATIN_FONT} (Latin) + {EA_FONT} (EA)")
    print(f"  Platform: {platform.system()}")
    print("-" * 60)

    filepath = generate_proposal(
        scenario_id=args.scenario,
        pricing_regime=args.pricing,
        year=args.year,
        output_path=args.output,
    )

    print("-" * 60)
    print(f"  PPT generated: {filepath}")
    print(f"  Slides: 20")
    print("=" * 60)


if __name__ == "__main__":
    main()
