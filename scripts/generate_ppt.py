#!/usr/bin/env python3
"""Generate Climate Risk 3-slide summary PPT.

Slides:
  1. Cover + Portfolio Summary
  2. Transition Risk
  3. Physical Risk

Usage:
    source backend/venv/bin/activate
    python generate_ppt.py
"""

import sys
import os
import platform
from datetime import date

# Add backend to path so we can import services
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "backend"))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from app.data.sample_facilities import get_all_facilities
from app.services.transition_risk import analyse_scenario, compare_scenarios
from app.services.physical_risk import assess_physical_risk
from app.core.config import SCENARIOS

# ── Color Palette ────────────────────────────────────────────────────
NAVY = RGBColor(0x1E, 0x29, 0x3B)
DARK_NAVY = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
ALT_ROW_GRAY = RGBColor(0xF8, 0xFA, 0xFC)
MID_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY = RGBColor(0x47, 0x55, 0x69)
BLACK = RGBColor(0x1E, 0x1E, 0x1E)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
YELLOW = RGBColor(0xEA, 0xB3, 0x08)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
LIGHT_GREEN = RGBColor(0xDC, 0xFC, 0xE7)
LIGHT_RED = RGBColor(0xFE, 0xE2, 0xE2)
LIGHT_ORANGE = RGBColor(0xFF, 0xED, 0xD5)
LIGHT_YELLOW = RGBColor(0xFE, 0xF9, 0xC3)

SCENARIO_COLORS = {
    "net_zero_2050": RED,
    "below_2c": ORANGE,
    "delayed_transition": YELLOW,
    "current_policies": GREEN,
}
SCENARIO_BG_COLORS = {
    "net_zero_2050": LIGHT_RED,
    "below_2c": LIGHT_ORANGE,
    "delayed_transition": LIGHT_YELLOW,
    "current_policies": LIGHT_GREEN,
}

# ── Slide dimensions (16:9) ─────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Layout constants ─────────────────────────────────────────────────
LEFT_MARGIN = Inches(0.8)
RIGHT_MARGIN = Inches(0.8)
CONTENT_WIDTH = Inches(11.7)

# ── Font strategy ────────────────────────────────────────────────────
_IS_MAC = platform.system() == "Darwin"
LATIN_FONT = "Calibri"
EA_FONT = "Apple SD Gothic Neo" if _IS_MAC else "Malgun Gothic"

# ── Sector name map ──────────────────────────────────────────────────
SECTOR_NAMES_KR = {
    "steel": "철강",
    "petrochemical": "석유화학",
    "automotive": "자동차",
    "electronics": "반도체/전자",
    "utilities": "발전",
    "cement": "시멘트",
    "shipping": "해운",
    "oil_gas": "정유",
}

HAZARD_NAMES_KR = {
    "flood": "홍수",
    "typhoon": "태풍",
    "heatwave": "폭염",
    "drought": "가뭄",
    "sea_level_rise": "해수면",
}


# ══════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ══════════════════════════════════════════════════════════════════════

def _set_font(run, name=LATIN_FONT, ea_name=EA_FONT, size=18, bold=False, color=BLACK):
    """Set font with EA (East Asian) fallback via lxml."""
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for ea in rPr.findall(qn('a:ea')):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', ea_name)


def _add_shadow(shape):
    """Add drop shadow to a shape via lxml."""
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    shadow = etree.SubElement(effectLst, qn('a:outerShdw'))
    shadow.set('blurRad', '50800')
    shadow.set('dist', '25400')
    shadow.set('dir', '5400000')
    shadow.set('algn', 'bl')
    shadow.set('rotWithShape', '0')
    srgb = etree.SubElement(shadow, qn('a:srgbClr'))
    srgb.set('val', '000000')
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', '20000')


def _set_slide_bg(slide, color):
    """Set a solid-color slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_gradient_bg(slide, color1_hex, color2_hex, angle=5400000):
    """Set gradient background on a slide via lxml."""
    bg = slide.background
    cSld = slide._element
    bg_elem = cSld.find(qn('p:bg'))
    if bg_elem is None:
        bg_elem = etree.SubElement(cSld, qn('p:bg'))
        cSld.insert(0, bg_elem)
    for old in bg_elem.findall(qn('p:bgPr')):
        bg_elem.remove(old)
    bgPr = etree.SubElement(bg_elem, qn('p:bgPr'))
    gradFill = etree.SubElement(bgPr, qn('a:gradFill'))
    gradFill.set('flip', 'none')
    gradFill.set('rotWithShape', '1')
    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))
    gs1 = etree.SubElement(gsLst, qn('a:gs'))
    gs1.set('pos', '0')
    srgb1 = etree.SubElement(gs1, qn('a:srgbClr'))
    srgb1.set('val', color1_hex)
    gs2 = etree.SubElement(gsLst, qn('a:gs'))
    gs2.set('pos', '100000')
    srgb2 = etree.SubElement(gs2, qn('a:srgbClr'))
    srgb2.set('val', color2_hex)
    lin = etree.SubElement(gradFill, qn('a:lin'))
    lin.set('ang', str(angle))
    lin.set('scaled', '1')
    etree.SubElement(bgPr, qn('a:effectLst'))


def _add_rounded_rect(slide, left, top, width, height, fill_color, shadow=True):
    """Add a rounded rectangle with optional drop shadow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if shadow:
        _add_shadow(shape)
    return shape


def _add_shape(slide, left, top, width, height, fill_color):
    """Add a plain rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a textbox with cross-platform font settings."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, size=font_size, color=color, bold=bold)
    return txBox


def _add_rich_text(slide, left, top, width, height, segments, alignment=PP_ALIGN.LEFT):
    """Add textbox with mixed formatting.
    segments = [(text, font_size, color, bold), ...]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    for text, font_size, color, bold in segments:
        run = p.add_run()
        run.text = text
        _set_font(run, size=font_size, color=color, bold=bold)
    return txBox


def _add_title_bar(slide, title_text, subtitle_text=None):
    """Add a navy title bar at the top of the slide."""
    _add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), NAVY)
    _add_textbox(slide, LEFT_MARGIN, Inches(0.15), Inches(10), Inches(0.6),
                 title_text, font_size=28, color=WHITE, bold=True)
    if subtitle_text:
        _add_textbox(slide, LEFT_MARGIN, Inches(0.7), Inches(10), Inches(0.4),
                     subtitle_text, font_size=14, color=MID_GRAY)
    _add_shape(slide, Inches(0), Inches(1.2), SLIDE_WIDTH, Inches(0.04), BLUE)


def _add_kpi_box(slide, left, top, width, height, label, value,
                 bg_color=LIGHT_BLUE, value_color=BLUE):
    """Add a rounded KPI box with shadow."""
    _add_rounded_rect(slide, left, top, width, height, bg_color, shadow=True)
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.4),
                 label, font_size=11, color=DARK_GRAY)
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.45), width - Inches(0.3), Inches(0.5),
                 value, font_size=28, color=value_color, bold=True)


def _add_table(slide, left, top, width, height, headers, rows, col_widths=None,
               header_bg=NAVY, header_fg=WHITE, stripe_color=ALT_ROW_GRAY):
    """Add a styled table to the slide."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        _style_cell(cell, font_size=12, bold=True, color=header_fg, bg_color=header_bg)

    for i, row in enumerate(rows):
        bg = stripe_color if i % 2 == 1 else WHITE
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            if isinstance(val, tuple):
                cell.text = val[0]
                _style_cell(cell, font_size=11, color=val[1],
                            bg_color=bg, bold=val[2] if len(val) > 2 else False)
            else:
                cell.text = str(val)
                _style_cell(cell, font_size=11, color=BLACK, bg_color=bg)

    return table_shape


def _style_cell(cell, font_size=11, bold=False, color=BLACK, bg_color=WHITE):
    """Style a table cell with cross-platform fonts."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg_color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            _set_font(run, size=font_size, bold=bold, color=color)
        if not p.runs:
            run = p.add_run()
            run.text = p.text
            _set_font(run, size=font_size, bold=bold, color=color)
        p.alignment = PP_ALIGN.CENTER


def _add_footer(slide, page_num):
    """Add consistent footer with page number."""
    _add_shape(slide, LEFT_MARGIN, Inches(7.15), CONTENT_WIDTH, Inches(0.02), MID_GRAY)
    _add_textbox(slide, Inches(12.0), Inches(7.18), Inches(1.0), Inches(0.3),
                 str(page_num), font_size=10, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)
    _add_textbox(slide, LEFT_MARGIN, Inches(7.18), Inches(6.0), Inches(0.3),
                 "Climate Risk Platform  |  Confidential", font_size=10, color=MID_GRAY)


def _fmt_b(val):
    """Format large USD values as $XXX.XB."""
    return f"${abs(val)/1e9:.1f}B"


def _fmt_m(val):
    """Format millions."""
    return f"{val/1e6:.1f}M"


def _fmt_pct(val):
    return f"{val:.0f}%"


# ══════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS (3 slides)
# ══════════════════════════════════════════════════════════════════════

def slide_1_cover(prs, facilities, scenario_name):
    """Slide 1: Cover + Portfolio Summary KPIs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, '0F172A', '1E293B')

    # Decorative accent
    _add_shape(slide, Inches(0), Inches(2.8), Inches(0.1), Inches(2.2), BLUE)

    # Title
    _add_textbox(slide, Inches(1.0), Inches(1.2), Inches(11), Inches(1.0),
                 "Climate Risk Platform", font_size=44, color=WHITE, bold=True)
    # Subtitle
    _add_textbox(slide, Inches(1.0), Inches(2.3), Inches(11), Inches(0.7),
                 f"기후 리스크 분석 보고서  |  {scenario_name}",
                 font_size=22, color=BLUE)
    # Date
    _add_textbox(slide, Inches(1.0), Inches(3.1), Inches(11), Inches(0.4),
                 f"{date.today().isoformat()}  |  Confidential",
                 font_size=14, color=MID_GRAY)

    # Portfolio KPIs
    total_s1 = sum(f["current_emissions_scope1"] for f in facilities)
    total_s2 = sum(f["current_emissions_scope2"] for f in facilities)
    total_assets = sum(f["assets_value"] for f in facilities)
    sectors = set(f["sector"] for f in facilities)

    kpis = [
        ("시설 수", f"{len(facilities)}개"),
        ("섹터 수", f"{len(sectors)}개"),
        ("총 배출량 (S1+S2)", f"{(total_s1+total_s2)/1e6:.0f}M tCO2e"),
        ("총 자산규모", f"${total_assets/1e9:.0f}B"),
    ]

    for i, (label, value) in enumerate(kpis):
        x = Inches(1.0) + i * Inches(3.1)
        y = Inches(4.2)
        _add_kpi_box(slide, x, y, Inches(2.8), Inches(1.1), label, value,
                     bg_color=NAVY, value_color=WHITE)

    # Bottom divider
    _add_shape(slide, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.02), MID_GRAY)
    _add_textbox(slide, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.4),
                 "NGFS 시나리오 기반 전환·물리적 리스크 통합 분석",
                 font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def slide_2_transition(prs, comparison, analysis):
    """Slide 2: Transition Risk — NPV cards + Top 5 facility table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "전환 리스크 분석",
                   "NGFS 시나리오별 탄소비용·에너지비용·매출영향 NPV")

    # NPV comparison cards
    npv_data = comparison["npv_comparison"]
    max_npv = max(abs(d["total_npv"]) for d in npv_data) if npv_data else 1

    for i, d in enumerate(npv_data):
        x = LEFT_MARGIN + i * Inches(3.05)
        y = Inches(1.6)
        sid = d["scenario"]
        color = SCENARIO_COLORS.get(sid, BLUE)
        bg = SCENARIO_BG_COLORS.get(sid, LIGHT_BLUE)

        _add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.0), bg, shadow=True)
        _add_shape(slide, x + Inches(0.1), y, Inches(2.6), Inches(0.06), color)

        _add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.4), Inches(0.4),
                     d["scenario_name"], font_size=14, color=NAVY, bold=True,
                     alignment=PP_ALIGN.CENTER)

        _add_textbox(slide, x + Inches(0.2), y + Inches(0.65), Inches(2.4), Inches(0.6),
                     f"-{_fmt_b(d['total_npv'])}", font_size=30, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)

        _add_textbox(slide, x + Inches(0.2), y + Inches(1.3), Inches(2.4), Inches(0.4),
                     f"Risk: {d['avg_risk_level']}", font_size=12, color=DARK_GRAY,
                     alignment=PP_ALIGN.CENTER)

        # Progress bar
        bar_pct = abs(d["total_npv"]) / max_npv if max_npv else 0
        bar_w = Emu(int(Inches(2.4) * bar_pct))
        _add_rounded_rect(slide, x + Inches(0.2), y + Inches(1.7),
                          bar_w, Inches(0.12), color, shadow=False)

    # Top 5 facility table
    _add_textbox(slide, LEFT_MARGIN, Inches(3.9), Inches(11), Inches(0.5),
                 "Top 5 고위험 시설 (Net Zero 2050)", font_size=18, color=NAVY, bold=True)

    risk_colors = {"High": RED, "Medium": ORANGE, "Low": GREEN}
    facs = sorted(analysis["facilities"], key=lambda x: x["delta_npv"])

    headers = ["순위", "시설명", "섹터", "Delta NPV", "자산 대비 %", "리스크"]
    rows = []
    for i, f in enumerate(facs[:5]):
        level = f["risk_level"]
        rows.append([
            str(i + 1),
            f["facility_name"],
            SECTOR_NAMES_KR.get(f["sector"], f["sector"]),
            f"-{_fmt_b(f['delta_npv'])}",
            f"{f['npv_as_pct_of_assets']:.1f}%",
            (level, risk_colors.get(level, BLACK), True),
        ])

    _add_table(slide, LEFT_MARGIN, Inches(4.4), CONTENT_WIDTH, Inches(2.5),
               headers, rows,
               col_widths=[Inches(0.8), Inches(3.0), Inches(2.0), Inches(2.0), Inches(1.8), Inches(2.1)])

    _add_footer(slide, 2)


def slide_3_physical(prs, physical_data):
    """Slide 3: Physical Risk — EAL summary + hazard cards + Top 5 table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "물리적 리스크 평가",
                   "5대 기후재해 유형별 시설 취약도 및 예상연간손실(EAL)")

    # 5 hazard mini-cards
    hazards_info = [
        ("홍수", "집중호우/하천범람", BLUE),
        ("태풍", "강풍/폭풍해일", RED),
        ("폭염", "냉각효율/생산성", ORANGE),
        ("가뭄", "용수부족/냉각수", YELLOW),
        ("해수면", "연안시설 침수", RGBColor(0x06, 0xB6, 0xD4)),
    ]

    for i, (title, desc, color) in enumerate(hazards_info):
        x = Inches(0.5) + i * Inches(2.5)
        y = Inches(1.5)
        _add_rounded_rect(slide, x, y, Inches(2.2), Inches(1.2), LIGHT_BLUE, shadow=True)
        _add_shape(slide, x, y, Inches(2.2), Inches(0.05), color)
        _add_textbox(slide, x + Inches(0.15), y + Inches(0.15), Inches(1.9), Inches(0.4),
                     title, font_size=16, color=NAVY, bold=True,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.15), y + Inches(0.55), Inches(1.9), Inches(0.5),
                     desc, font_size=11, color=DARK_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # EAL Top 5 table
    _add_textbox(slide, LEFT_MARGIN, Inches(2.9), Inches(11), Inches(0.5),
                 "시설별 연간예상손실(EAL) Top 5", font_size=18, color=NAVY, bold=True)

    fac_results = sorted(physical_data["facilities"],
                         key=lambda x: x["total_expected_annual_loss"], reverse=True)

    headers = ["순위", "시설명", "위치", "위험등급", "EAL ($M)", "주요 재해"]
    rows = []
    for i, fac in enumerate(fac_results[:5]):
        top_hazard = max(fac["hazards"], key=lambda h: h["potential_loss"])
        level = fac["overall_risk_level"]
        rows.append([
            str(i + 1),
            fac["facility_name"],
            fac["location"],
            (level,
             {"High": RED, "Medium": ORANGE, "Low": GREEN}.get(level, BLACK),
             True),
            f"${fac['total_expected_annual_loss']/1e6:.1f}",
            HAZARD_NAMES_KR.get(top_hazard["hazard_type"], top_hazard["hazard_type"]),
        ])

    _add_table(slide, LEFT_MARGIN, Inches(3.4), CONTENT_WIDTH, Inches(2.0),
               headers, rows,
               col_widths=[Inches(0.8), Inches(2.5), Inches(2.0), Inches(1.8), Inches(2.2), Inches(2.4)])

    # Risk distribution summary bar
    summary = physical_data["overall_risk_summary"]
    total_eal = sum(f["total_expected_annual_loss"] for f in physical_data["facilities"])

    _add_rounded_rect(slide, LEFT_MARGIN, Inches(5.7), CONTENT_WIDTH, Inches(0.9), NAVY, shadow=True)

    _add_rich_text(slide, Inches(1.0), Inches(5.8), Inches(5.0), Inches(0.5),
                   [
                       ("총 EAL: ", 14, WHITE, True),
                       (f"${total_eal/1e6:.1f}M", 14, LIGHT_BLUE, True),
                       ("    ", 14, WHITE, False),
                       (f"High {summary.get('High', 0)}개", 14, LIGHT_RED, True),
                       ("  |  ", 14, MID_GRAY, False),
                       (f"Medium {summary.get('Medium', 0)}개", 14, LIGHT_ORANGE, True),
                       ("  |  ", 14, MID_GRAY, False),
                       (f"Low {summary.get('Low', 0)}개", 14, LIGHT_GREEN, True),
                   ])

    _add_textbox(slide, Inches(7.0), Inches(5.85), Inches(5.5), Inches(0.4),
                 "Analytical v1 | IPCC AR6 + KMA 30yr | Open-Meteo API",
                 font_size=12, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

    _add_footer(slide, 3)


# ══════════════════════════════════════════════════════════════════════
# ENTRY POINTS
# ══════════════════════════════════════════════════════════════════════

def generate_report(
    scenario_id: str = "net_zero_2050",
    pricing_regime: str = "kets",
    year: int = 2030,
) -> str:
    """Generate 3-slide PPT and return file path."""
    facilities = get_all_facilities()
    analysis = analyse_scenario(scenario_id, pricing_regime=pricing_regime)
    comparison = compare_scenarios(pricing_regime=pricing_regime)
    physical = assess_physical_risk(scenario_id=scenario_id, year=year)

    scenario_name = SCENARIOS[scenario_id]["name"]

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_1_cover(prs, facilities, scenario_name)
    slide_2_transition(prs, comparison, analysis)
    slide_3_physical(prs, physical)

    output_dir = os.path.join(os.path.dirname(__file__), "outputs")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "climate_risk_report.pptx")
    prs.save(output_path)
    return output_path


def main():
    print("=" * 60)
    print("  Climate Risk Platform - PPT Report (3 slides)")
    print("=" * 60)

    filepath = generate_report()
    print(f"\nPPT generated: {filepath}")
    print(f"  Font: {LATIN_FONT} (Latin) + {EA_FONT} (EA)")
    print(f"  Platform: {platform.system()}")
    print("=" * 60)


if __name__ == "__main__":
    main()
